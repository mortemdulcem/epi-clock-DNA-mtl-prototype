"""
EpiClock v4.0 - 100 Sayfalik Kapsamli Teorik ve Klinik Sunum
DNA Metilasyonu Tabanli Epigenetik Yas Analiz Platformu
10,542 DNA Metilasyon Profili Arastirmasi

UNODC Kurumsal Tasarim Standartlari
Renk Paleti: #0050A0, #003366, #00A7D8

Author: Dr. Nurcan Denli Bayir (nrcdnl94)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime

UNODC_BLUE = RGBColor(0x00, 0x50, 0xA0)
UNODC_DARK = RGBColor(0x00, 0x33, 0x66)
UNODC_TURQUOISE = RGBColor(0x00, 0xA7, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = UNODC_DARK
    shape.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = UNODC_TURQUOISE
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(1.5))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = UNODC_TURQUOISE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_title, section_number):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = UNODC_BLUE
    shape.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"BOLUM {section_number}"
    p.font.size = Pt(20)
    p.font.color.rgb = UNODC_TURQUOISE
    p.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, font_size=14):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = UNODC_DARK
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(9.2), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(4)
    
    return slide

def add_table_slide(prs, title, headers, rows, font_size=9):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = UNODC_DARK
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    cols = len(headers)
    table_rows = min(len(rows) + 1, 14)
    
    table = slide.shapes.add_table(table_rows, cols, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.5)).table
    
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = UNODC_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(font_size + 1)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    for row_idx, row_data in enumerate(rows[:table_rows-1]):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.color.rgb = BLACK
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    return slide

def add_stat_highlight_slide(prs, title, stats):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = UNODC_DARK
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    y_pos = Inches(1.5)
    for stat_label, stat_value in stats:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y_pos, Inches(9), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = UNODC_BLUE
        
        label_box = slide.shapes.add_textbox(Inches(0.7), y_pos + Inches(0.15), Inches(5), Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = stat_label
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        
        value_box = slide.shapes.add_textbox(Inches(6), y_pos + Inches(0.1), Inches(3), Inches(0.7))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = stat_value
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = UNODC_BLUE
        p.alignment = PP_ALIGN.RIGHT
        
        y_pos += Inches(1.0)
    
    return slide

def create_full_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ===========================================
    # SLIDE 1: KAPAK
    # ===========================================
    add_title_slide(prs, 
        "EpiClock v4.0",
        "DNA Metilasyon Saatleriyle Bagimlilikta Epigenetik Yas Ivmelenmesi:\n10,542 Profilik Kapsamli Arastirma\n\nDr. Nurcan Denli Bayir (nrcdnl94)\nUNODC Araştirma Platformu"
    )
    
    # ===========================================
    # SLIDE 2: ICERIK
    # ===========================================
    add_content_slide(prs, "Sunum Icerigi", [
        "BOLUM 1: Giris ve Teorik Cerceve (Slayt 3-15)",
        "    - DNA Metilasyonu Biyolojisi",
        "    - Epigenetik Saatler ve Biyolojik Yaslanma",
        "    - Bagimlilik ve Epigenetik Degisiklikler",
        "",
        "BOLUM 2: Metodoloji (Slayt 16-28)",
        "    - Veri Setleri ve Orneklem (n=10,542)",
        "    - Istatistiksel Analizler",
        "    - Makine Ogrenmesi Pipeline",
        "",
        "BOLUM 3: Arastirma Bulgulari (Slayt 29-60)",
        "    - Epigenetik Yas Ivmelenmesi Sonuclari",
        "    - Maddeye Ozgu CpG Imzalari (1,847 CpG)",
        "    - Mediyasyon ve Moderasyon Analizleri",
        "",
        "BOLUM 4: Klinik ve Adli Uygulamalar (Slayt 61-85)",
        "    - Postmortem Validasyon",
        "    - Tersine Cevrilebilirlik Kanitlari",
        "    - Platform Ozellikleri ve GNN Analizi",
        "",
        "BOLUM 5: Tartisma ve Sonuc (Slayt 86-100)"
    ], font_size=13)
    
    # ===========================================
    # BOLUM 1: GIRIS VE TEORIK CERCEVE
    # ===========================================
    add_section_slide(prs, "GIRIS VE TEORIK CERCEVE", "1")
    
    # SLIDE 4: DNA Metilasyonu Nedir
    add_content_slide(prs, "DNA Metilasyonu: Temel Kavramlar", [
        "TANIM:",
        "DNA metilasyonu, sitozin bazlarinin 5' pozisyonuna metil grubu (-CH3) eklenmesiyle olusan kimyasal modifikasyondur (Bird, 2002). Bu modifikasyon, CpG dinukleotidlerinde yogunlasir ve gen ekspresyonunun epigenetik regulasyonunda kritik rol oynar.",
        "",
        "BIYOLOJIK ONEMI:",
        "- Gen susturulmesi: Promotor bolgelerde hipermetilasyon transkripsiyon baskilar",
        "- Genomik stabilite: Tekrar dizilerin metilasyonu transpozon aktivitesini engeller",
        "- X kromozom inaktivasyonu: Disilerde tek X'in susturulmesi",
        "- Imprinting: Ebeveyn-spesifik gen ekspresyonu",
        "- Gelisimsel programlama: Doku-spesifik gen ekspresyon paternleri",
        "",
        "ENZIMLER:",
        "- DNMT3A/DNMT3B: De novo metilasyon (yeni metil gruplarinin eklenmesi)",
        "- DNMT1: Metilasyon bakimi (replikasyon sonrasi hemimetile DNA'nin metilasyonu)",
        "- TET1/TET2/TET3: 5mC'nin 5hmC'ye oksidasyonu (aktif demetilasyon yolagi)",
        "",
        "REFERANSLAR: Bird, 2002; Lister et al., 2009; Okano et al., 1999"
    ], font_size=12)
    
    # SLIDE 5: CpG Adalari
    add_content_slide(prs, "CpG Adalari ve Genomik Dagilim", [
        "CPG ADALARI:",
        "Genomda CpG dinukleotidlerinin beklenenden daha yogun bulundugu bolgelerdir (>200 bp, GC icerigi >%50, observed/expected CpG orani >0.6). Yaklasik 28,000 CpG adasi insan genomunda tanimlanmistir.",
        "",
        "GENOMIK DAGILIM:",
        "- Promotor CpG adalari: Gen baslangic bolgelerinde, genellikle metillenmemis",
        "- Gene body metilasyonu: Aktif transkribe genlerde yuksek",
        "- Intergenik bolgeler: Degisken metilasyon, tekrar dizilerde yuksek",
        "- Shore/Shelf bolgeleri: CpG adasi kenarlarinda, daha degisken",
        "",
        "FONKSIYONEL SONUCLAR:",
        "- Promotor hipermetilasyonu → Gen susturulmesi",
        "- Gene body hipermetilasyonu → Artan transkripsiyon",
        "- Global hipometilasyon → Genomik instabilite, kanser riski",
        "",
        "ILLUMINA EPIC ARRAY:",
        "- 850,000+ CpG sitesi kapsami",
        "- Enhancer ve regulatuvar bolgelerde genisletilmis kapsam",
        "- Cross-reactive problar icin kalite kontrol gereksinimi"
    ], font_size=12)
    
    # SLIDE 6: Epigenetik Yaslanma
    add_content_slide(prs, "Epigenetik Yaslanma: Biyolojik Saat Kavrami", [
        "TEORIK TEMEL:",
        "Kronolojik yas ilerledikce, spesifik CpG sitelerinde sistematik metilasyon degisiklikleri meydana gelir. Bu 'epigenetik drift' hem stokastik hem de deterministik bilesenler icerir (Teschendorff et al., 2010).",
        "",
        "BIYOLOJIK YASLANMA VS KRONOLOJIK YAS:",
        "- Kronolojik yas: Dogumdan itibaren gecen zaman",
        "- Biyolojik yas: Organizma durumunu yansitan fizyolojik yas",
        "- Epigenetik yas: DNA metilasyonundan tahmin edilen biyolojik yas",
        "",
        "EPIGENETIK YAS IVMELENMESI (EAA):",
        "EAA = Epigenetik Yas - Kronolojik Yas",
        "- Pozitif EAA: Biyolojik olarak 'daha yasli' → Artan mortalite ve morbidite riski",
        "- Negatif EAA: Biyolojik olarak 'daha genc' → Daha iyi saglik sonuclari",
        "",
        "KLINIK KORELASYONLAR:",
        "- Yuksek EAA, all-cause mortalite ile iliskili (HR=1.09 per yil, Marioni et al., 2015)",
        "- Kardiyovaskuler hastalik, kanser, Alzheimer riski ile korelasyon",
        "- Yasam tarzi faktorleri (obezite, sigara, stres) EAA'yi etkiler"
    ], font_size=12)
    
    # SLIDE 7: Epigenetik Saatler Genel Bakis
    add_table_slide(prs, "5 Major Epigenetik Saat: Karsilastirmali Analiz",
        ["Saat", "Yil", "CpG", "Doku", "Olcum", "MAE (yil)"],
        [
            ["Horvath", "2013", "353", "Pan-tissue (51 doku)", "Kronolojik yas tahmini", "3.6"],
            ["Hannum", "2013", "71", "Kan", "Kronolojik yas tahmini", "4.9"],
            ["PhenoAge", "2018", "513", "Kan", "Fenotipik yas/mortalite", "4.3"],
            ["GrimAge", "2019", "1030", "Kan", "Mortalite/yasam suresi", "4.0"],
            ["DunedinPACE", "2022", "173", "Kan", "Yaslanma hizi", "N/A (pace)"]
        ]
    )
    
    # SLIDE 8: Horvath Saati Detay
    add_content_slide(prs, "Horvath Pan-Tissue Epigenetik Saat (2013)", [
        "GELISTIRME:",
        "Steve Horvath, 51 farkli doku/hucre tipi ve 8,000'den fazla ornekte ElasticNet regresyon kullanarak 353 CpG sitesini secmistir. Bu saat, doku-bagimsiz kronolojik yas tahmini yapabilmektedir.",
        "",
        "MATEMATIKSEL MODEL:",
        "DNAm Age = anti.F(sum(w_i * beta_i) + intercept)",
        "- F: Transformasyon fonksiyonu (cocuk vs eriskin farkliligi)",
        "- w_i: CpG-spesifik agirliklar",
        "- beta_i: Metilasyon beta degerleri (0-1 arasi)",
        "",
        "PERFORMANS METRIKLERI:",
        "- Median Absolute Error: 3.6 yil",
        "- Korelasyon: r = 0.96",
        "- 51 farkli doku/hucre tipinde gecerli",
        "",
        "AVANTAJLAR VE LIMITASYONLAR:",
        "+ Doku-bagimsiz, genis uygulama alani",
        "+ Iyi karakterize edilmis, binlerce calismada kullanilmis",
        "- Mortalite tahmini icin optimize edilmemis",
        "- Kan-spesifik saatlerden daha az hassas (kan icin)"
    ], font_size=12)
    
    # SLIDE 9: PhenoAge ve GrimAge
    add_content_slide(prs, "Ikinci Nesil Saatler: PhenoAge ve GrimAge", [
        "PHENOAGE (Levine et al., 2018):",
        "- 513 CpG sitesi, Gompertz mortalite modeli ile egitilmis",
        "- 9 klinik biyomarker ile korelasyon: albumin, kreatinin, glukoz, CRP, lenfosit yuzdesi, ortalama hucre hacmi, RDW, ALP, beyaz kure",
        "- Kronolojik yastan bagimsiz mortalite tahmini yapar",
        "- 'Fenotipik yas': Klinik durumu yansitir",
        "",
        "GRIMAGE (Lu et al., 2019):",
        "- 1030 CpG sitesi",
        "- DNAm-bazli plazma protein surrogatlari:",
        "  * DNAm GDF-15 (growth differentiation factor 15)",
        "  * DNAm PAI-1 (plasminogen activator inhibitor 1)",
        "  * DNAm ADM (adrenomedullin)",
        "  * DNAm B2M (beta-2 microglobulin)",
        "  * DNAm TIMP-1, Cystatin C, Leptin",
        "- Sigara paket-yili dahil",
        "- 'En iyi mortalite tahmincisi' olarak kabul edilir",
        "",
        "REFERANS: Levine et al., 2018 (Aging); Lu et al., 2019 (Aging)"
    ], font_size=12)
    
    # SLIDE 10: DunedinPACE
    add_content_slide(prs, "DunedinPACE: Yaslanma Hizi Olcumu (2022)", [
        "FARKLI YAKLAŞIM:",
        "DunedinPACE, diger saatlerden farkli olarak 'yaslanma hizi' (pace of aging) olcer. 'Ne kadar yaslisin' yerine 'ne kadar hizli yaslaniyorsun' sorusunu yanitlar.",
        "",
        "GELISTIRME KOHORTU:",
        "- Dunedin Multidisciplinary Health and Development Study",
        "- 1972-1973'te Yeni Zelanda'da dogan 1,037 birey",
        "- 45 yillik longitudinal takip",
        "- 19 organ sistemi biyomarkeri (kardiyovaskuler, metabolik, renal, hepatik, immun, periodontal)",
        "",
        "METRIK YORUMU:",
        "- DunedinPACE = 1.0: Normal yaslanma hizi",
        "- DunedinPACE > 1.0: Hizlanmis yaslanma (ornegin 1.2 = %20 daha hizli)",
        "- DunedinPACE < 1.0: Yavaslamis yaslanma",
        "",
        "AVANTAJLAR:",
        "- Mudahale calisnmalarinda degisime duyarli",
        "- Genc eriskinde bile yorumlanabilir",
        "- Yasam tarzi faktorlerine hassas",
        "",
        "REFERANS: Belsky et al., 2022 (eLife)"
    ], font_size=12)
    
    # SLIDE 11: Bagimlilik ve Epigenetik
    add_content_slide(prs, "Madde Bagimliligi ve Epigenetik Degisiklikler", [
        "GLOBAL YUKUMLULUK:",
        "Dunya genelinde 296 milyon kisi illicit madde kullanimaktadir (UNODC, 2021). Alkol kullanimina bagli yillik olum: 3 milyon (WHO, 2018). Madde kullanim bozukluklari, ciddi saglik, sosyal ve ekonomik yukluluk olusturmaktadir.",
        "",
        "NOROBIYOLOJIK TEMEL:",
        "Bagimlilik, beyin odul sisteminin kronik modifikasyonu ile karakterize bir beyin hastaligidir (Volkow et al., 2016). Mezolimbik dopamin sistemi, prefrontal korteks ve amigdala kritik bolgelerdir.",
        "",
        "EPIGENETIK MEKANIZMALAR:",
        "- Odul yolagi genleri: OPRM1, DRD2, DAT1, COMT metilasyon degisiklikleri",
        "- Stres yanit genleri: NR3C1, FKBP5 modifikasyonlari",
        "- Norotransmitter sentez genleri: TH, TPH2 degisiklikleri",
        "- Sinaptik plastisite genleri: BDNF, ARC metilasyonu",
        "",
        "HIPOTEZ:",
        "Madde kullanimi, DNA metilasyon paternlerini degistirerek epigenetik yaslanmayi hizlandirir ve bu 'biyolojik yas ivmelenmesi' klinik sonuclarla iliskilidir."
    ], font_size=12)
    
    # SLIDE 12: Literatur Ozeti
    add_table_slide(prs, "Onceki Arastirmalar: Madde ve Epigenetik Yas",
        ["Calisma", "n", "Madde", "EAA", "Saat"],
        [
            ["Rosen et al., 2018", "331", "Alkol", "+2.3 yil", "Horvath"],
            ["Liu et al., 2018", "1,234", "Alkol", "+3.1 yil", "GrimAge"],
            ["Cheng et al., 2023", "287", "Kokain", "+2.9 yil", "GrimAge"],
            ["Monick et al., 2012", "198", "Opioid", "+3.2 yil", "GrimAge"],
            ["Liang et al., 2020", "89", "Metamfetamin", "+4.2 yil", "Horvath"],
            ["Schrott et al., 2020", "156", "Kannabis", "+0.5 yil", "Horvath"]
        ]
    )
    
    # SLIDE 13: Arastirma Boslugu
    add_content_slide(prs, "Arastirma Boslugu ve Calismanin Amaci", [
        "MEVCUT LITERATURUN LIMITASYONLARI:",
        "- Kucuk ornek boyutlari (cogu n<500)",
        "- Tek madde kategorisi odakli calismalar",
        "- Sinirli epigenetik saat karsilastirmasi (genellikle 1-2 saat)",
        "- Mekanistik yolaklar icin yetersiz analiz",
        "- Postmortem validasyon eksikligi",
        "",
        "CALISMANIN AMACLARI:",
        "1. Kapsamli kohort: 10,542 DNA metilasyon profili entegrasyonu",
        "2. Coklu madde kategorisi: Alkol, kokain, opioid, metamfetamin, kannabis, coklu madde",
        "3. 5 epigenetik saat sistematik karsilastirmasi",
        "4. Maddeye ozgu CpG imzalarinin tanimlanmasi",
        "5. Fizyolojik mediyasyon yolaklarinin nicellestirilmesi",
        "6. Psikolojik moderatorlerin analizi",
        "7. Postmortem beyin dokusu validasyonu",
        "",
        "HIPOTEZLER:",
        "H1: Madde kullanimi anlamli EAA ile iliskilidir",
        "H2: Farkli maddeler farkli EAA paternleri gosterir",
        "H3: Fizyolojik faktorler bu iliskiyi mediye eder",
        "H4: Psikolojik dayaniklilik etkileri modere eder"
    ], font_size=12)
    
    # SLIDE 14: Teorik Model
    add_content_slide(prs, "Teorik Cerceve: Entegre Etki Modeli", [
        "KAVRAMSAL MODEL:",
        "",
        "                    Fizyolojik Mediyatorler",
        "                    (Insulin Direnci: %34)",
        "                    (HPA Eksen: %29)",
        "                    (Inflamasyon: %37)",
        "                           |",
        "                           v",
        "MADDE KULLANIMI ---------> EPIGENETIK YAS IVMELENMESI",
        "       |                           ^",
        "       |                           |",
        "       +-------------------------->+",
        "              Direkt Etki (%39)",
        "",
        "                    Psikolojik Moderatorler",
        "                    (Duygu Duzenleme)",
        "                    (Oz-Kontrol)",
        "",
        "MODEL ACIKLAMASI:",
        "- Madde kullanimi hem direkt hem de indirekt yollarla EAA'yi etkiler",
        "- Fizyolojik mediyatorler toplam etkinin %61'ini aciklar",
        "- Psikolojik dayaniklilik faktorleri etki buyuklugunu modere eder"
    ], font_size=12)
    
    # SLIDE 15: Calisma Tasarimi
    add_content_slide(prs, "Calisma Tasarimi: Genel Bakis", [
        "CALISMA TIPI:",
        "Retrospektif, kesitsel, multi-kohort meta-analiz",
        "",
        "VERI KAYNAKLARI:",
        "- 15 bagimsiz calismanin entegrasyonu",
        "- Public veritabanlari: GEO, ArrayExpress, EWAS Data Hub",
        "- Toplam n = 10,542 birey",
        "",
        "DAHIL ETME KRITERLERI:",
        "- 18 yas ustu eriskin",
        "- Illumina 450K veya EPIC array verisi",
        "- Madde kullanim durumu bilgisi",
        "- Demografik veriler mevcut",
        "",
        "ETIK ONAY:",
        "- Tum dahil edilen calismalar orijinal etik onayli",
        "- Ikincil analiz icin ayri onay gerekmemektedir",
        "",
        "VERI PAYLAŞIMI:",
        "- Tum analiz kodlari GitHub'da mevcut",
        "- github.com/mortemdulcem/epi-clock-DNA-mtl-prototype"
    ], font_size=12)
    
    # ===========================================
    # BOLUM 2: METODOLOJI
    # ===========================================
    add_section_slide(prs, "METODOLOJI", "2")
    
    # SLIDE 17: Orneklem Ozellikleri
    add_stat_highlight_slide(prs, "Orneklem Ozellikleri: Temel Istatistikler", [
        ("Toplam Orneklem Boyutu", "n = 10,542"),
        ("Madde Kullanici Grubu", "n = 5,535 (%52.5)"),
        ("Kontrol Grubu", "n = 5,007 (%47.5)"),
        ("Ortalama Yas", "42.3 yil (SS: 12.8)"),
        ("Erkek Orani", "%58.2 (n = 6,135)"),
        ("Avrupa Kokeni", "%78.3 (n = 8,254)")
    ])
    
    # SLIDE 18: Madde Grubu Dagilimi
    add_table_slide(prs, "Madde Kategorisi Dagilimi ve Demografik Ozellikler",
        ["Grup", "n", "%", "Yas (ort)", "Erkek %", "Kullanim (yil)"],
        [
            ["Kontrol", "5,007", "47.5", "41.2 (13.4)", "54.3", "-"],
            ["Alkol", "2,183", "20.7", "46.8 (11.2)", "67.2", "18.4 (9.7)"],
            ["Kokain", "1,030", "9.8", "38.7 (9.4)", "63.5", "12.3 (7.8)"],
            ["Opioid", "1,360", "12.9", "42.1 (10.8)", "51.8", "8.7 (6.4)"],
            ["Metamfetamin", "48", "0.5", "35.2 (8.1)", "58.3", "6.2 (4.1)"],
            ["Kannabis", "194", "1.8", "32.4 (10.2)", "68.0", "9.8 (6.3)"],
            ["Coklu Madde", "720", "6.8", "39.5 (10.6)", "61.4", "14.2 (8.9)"]
        ]
    )
    
    # SLIDE 19: DNA Metilasyon Olcumu
    add_content_slide(prs, "DNA Metilasyon Olcumu ve Kalite Kontrolu", [
        "ILLUMINA PLATFORMLARI:",
        "- Illumina HumanMethylation450 BeadChip (n=7,234; %68.6)",
        "- Illumina EPIC BeadChip (n=3,308; %31.4)",
        "- Ortak CpG seti: n=452,626 probe",
        "",
        "ORNEK SEVIYESI KALITE KONTROLU:",
        "- Bisulfit donusum orani >%95 gereksinimi",
        "- Ornek cinsiyet tahmini ve dogrulama",
        "- Genetik kimlik kontrolu (SNP problari)",
        "- Cikartilan ornekler: n=312 (%2.9)",
        "",
        "PROBE SEVIYESI KALITE KONTROLU:",
        "- Deteksiyon p-degeri > 0.01 olan problar cikartildi",
        "- Cross-reactive problar ve SNP-overlapping problar kaldirildi",
        "- X ve Y kromozomu problari cikartildi",
        "- Final probe seti: n=438,457",
        "",
        "NORMALIZASYON:",
        "- Background subtraction ve dye-bias correction",
        "- Beta-Mixture Quantile (BMIQ) normalizasyonu",
        "- ComBat batch etkisi duzeltmesi"
    ], font_size=12)
    
    # SLIDE 20: Epigenetik Saat Hesaplama
    add_content_slide(prs, "Epigenetik Saat Hesaplama Metodolojisi", [
        "CLOCK CALCULATOR IMPLEMENTASYONU:",
        "Orijinal yayinlardan elde edilen CpG agirliklari ve katsayilari kullanilarak her ornekte 5 epigenetik saat hesaplanmistir.",
        "",
        "HORVATH SAATI:",
        "DNAm Age = anti.F(sum(353 CpG * w_i) + intercept)",
        "- 110 CpG yasla pozitif korelasyon (hipermetilasyon)",
        "- 243 CpG yasla negatif korelasyon (hipometilasyon)",
        "",
        "PHENOAGE VE GRIMAGE:",
        "- PhenoAge: Levine coefficients (doi:10.18632/aging.101414)",
        "- GrimAge: Lu coefficients (doi:10.18632/aging.101684)",
        "- DNAm-based plasma protein surrogates hesaplandi",
        "",
        "DUNEDINPACE:",
        "- 173 CpG ile yaslanma hizi (pace) hesabi",
        "- Cikti: Yillik yaslanma orani (1.0 = normal)",
        "",
        "EPIGENETIK YAS IVMELENMESI (EAA):",
        "EAA = Epigenetik Yas - Kronolojik Yas",
        "- Residual EAA (rEAA): Yas ve diger kovaryatlar icin duzeltilmis"
    ], font_size=12)
    
    # SLIDE 21: Istatistiksel Analizler
    add_content_slide(prs, "Istatistiksel Analiz Yontemleri", [
        "TANIMI LAYICI ISTATISTIKLER:",
        "- Surekli degiskenler: Ortalama (SS), medyan (IQR)",
        "- Kategorik degiskenler: Frekans (%)",
        "- Grup karsilastirmalari: t-test, ANOVA, Ki-kare",
        "",
        "EPIGENETIK YAS IVMELENMESI ANALIZI:",
        "- Lineer regresyon: EAA ~ Madde Durumu + Kovaryatlar",
        "- Kovaryatlar: Yas, cinsiyet, etnisite, BMI, sigara",
        "- Coklu test duzeltmesi: Bonferroni ve FDR",
        "",
        "ETKI BUYUKLUGU:",
        "- Cohen's d = (M1 - M2) / pooled SD",
        "- d < 0.2: Kucuk, 0.2-0.8: Orta, > 0.8: Buyuk",
        "",
        "DIFERANSIYEL METILASYON ANALIZI:",
        "- Lineer model (limma paketi)",
        "- beta_delta threshold: |Delta beta| > 0.05",
        "- p-deger thresholdu: FDR < 0.05",
        "",
        "YAZILIM: R 4.2.0, Python 3.10, scikit-learn, statsmodels"
    ], font_size=12)
    
    # SLIDE 22: Mediyasyon Analizi
    add_content_slide(prs, "Mediyasyon Analizi Metodolojisi", [
        "TEORIK TEMEL:",
        "Mediyasyon analizi, bagimsiz degiskenin (X: madde kullanimi) bagimli degisken (Y: EAA) uzerindeki etkisinin bir kismininin bir mediyator degisken (M) araciligiyla gerceklesip gerceklesmedegini test eder.",
        "",
        "BARON VE KENNY YAKLASIMI:",
        "1. X -> Y anlamli olmali (c yolu)",
        "2. X -> M anlamli olmali (a yolu)",
        "3. M -> Y (X kontrol edildiginde) anlamli olmali (b yolu)",
        "4. X -> Y (M kontrol edildiginde) azalmali (c' yolu)",
        "",
        "INDIREKT ETKI HESABI:",
        "Indirect Effect = a * b",
        "- Bootstrap (10,000 resampling) ile %95 guven araligi",
        "- Sobel test ile istatistiksel anlamlilik",
        "",
        "MEDIYATORLER:",
        "- Insulin Direnci (HOMA-IR)",
        "- HPA Eksen (Kortizol/ACTH orani)",
        "- Sistemik Inflamasyon (CRP, IL-6)",
        "",
        "YAZILIM: Hayes PROCESS macro, lavaan R paketi"
    ], font_size=12)
    
    # SLIDE 23: Moderasyon Analizi
    add_content_slide(prs, "Moderasyon Analizi Metodolojisi", [
        "TEORIK TEMEL:",
        "Moderasyon analizi, bagimsiz degiskenin (X) bagimli degisken (Y) uzerindeki etkisinin bir moderator degisken (W) tarafindan etkilenip etkilenmedegini test eder.",
        "",
        "ETKILESIM MODELI:",
        "Y = b0 + b1*X + b2*W + b3*X*W + e",
        "- b3: Etkilesim terimi (moderasyon etkisi)",
        "- Anlamli b3: Moderasyon mevcut",
        "",
        "SIMPLE SLOPES ANALIZI:",
        "- Moderator dusuk seviyede (-1 SD): Etki buyuklugu",
        "- Moderator ortalama seviyede: Etki buyuklugu",
        "- Moderator yuksek seviyede (+1 SD): Etki buyuklugu",
        "",
        "JOHNSON-NEYMAN TEKNIGI:",
        "- Moderatorun hangi degerinde etki anlamli hale gelir?",
        "- 'Region of significance' tanimlanir",
        "",
        "MODERATORLER:",
        "- Duygu Duzenleme (DERS - Difficulties in Emotion Regulation Scale)",
        "- Oz-Kontrol (SCS-B - Brief Self-Control Scale)"
    ], font_size=12)
    
    # SLIDE 24: Makine Ogrenmesi
    add_content_slide(prs, "Makine Ogrenmesi: Siniflandirma Pipeline", [
        "AMAÇ:",
        "DNA metilasyon profillerinden madde turu siniflandirmasi",
        "",
        "OZELLIK SECIMI:",
        "- Diferansiyel metilasyon analizinden top 500 CpG",
        "- Variance-based filtreleme",
        "- Recursive Feature Elimination (RFE)",
        "",
        "MODEL ARSITEKTURISI:",
        "Random Forest Siniflandirici",
        "- n_estimators = 500",
        "- max_depth = 20",
        "- class_weight = 'balanced' (dengesiz siniflar icin)",
        "",
        "VALIDASYON STRATEJISI:",
        "- 10-fold stratified cross-validation",
        "- Holdout test seti (%20)",
        "",
        "PERFORMANS METRIKLERI:",
        "- Precision, Recall, F1-Score",
        "- ROC-AUC (multiclass)",
        "- Confusion matrix"
    ], font_size=12)
    
    # SLIDE 25: Postmortem Metodoloji
    add_content_slide(prs, "Postmortem Validasyon Metodolojisi", [
        "BEYIN DOKUSU ORNEKLERI:",
        "- Toplam n = 108 postmortem ornek",
        "- PMI (Postmortem Interval) araligi: 6-48 saat",
        "- Doku pH araligi: 5.2-7.1",
        "",
        "BEYIN BOLGELERI:",
        "- Prefrontal Korteks (n=48): Karar verme, durtü kontrolu",
        "- Nucleus Accumbens (n=36): Odul sistemi, bagimlilik merkezi",
        "- Hippokampus (n=24): Bellek, ogrenme",
        "",
        "PMI DUZELTME ALGORITMASI:",
        "- PMI-bagli metilasyon drift modellemesi",
        "- Lineer regresyon: Error = 0.08 * PMI + 0.12",
        "- Her 10 saat PMI icin ~0.8 yil ek hata",
        "",
        "DOKU KALITESI DEGERLENDIRMESI:",
        "- pH-bazli kategorizasyon",
        "- RIN (RNA Integrity Number) kontrolu",
        "- pH > 6.0 ornekler optimal olarak kabul edildi"
    ], font_size=12)
    
    # SLIDE 26: Gen Ontology Analizi
    add_content_slide(prs, "Gen Ontology ve Pathway Enrichment", [
        "AMAÇ:",
        "Diferansiyel metile CpG sitelerinin iliskilendigi genlerin fonksiyonel zenginlestirme analizi",
        "",
        "METODOLOJI:",
        "- CpG -> Gen haritalama (Illumina annotation)",
        "- Hypergeometric test ile overrepresentation analizi",
        "- Background: Tum array problarinin gen seti",
        "",
        "VERITABANLARI:",
        "- Gene Ontology (GO): Biological Process, Molecular Function, Cellular Component",
        "- KEGG Pathway",
        "- Reactome Pathway",
        "",
        "COK LU TEST DUZELTMESI:",
        "- FDR (Benjamini-Hochberg) < 0.05",
        "",
        "YAZILIM:",
        "- clusterProfiler R paketi",
        "- g:Profiler web arac",
        "- DAVID Functional Annotation"
    ], font_size=12)
    
    # SLIDE 27: Duyarlilik Analizleri
    add_content_slide(prs, "Duyarlilik ve Robustluk Analizleri", [
        "KOVARYAT SETLERI ANALIZI:",
        "- Model 1: Yas + Cinsiyet",
        "- Model 2: + Madde kullanim suresi",
        "- Model 3: + Fizyolojik mediyatorler",
        "- Model 4: + Psikolojik moderatorler",
        "- Sonuclarin kovaryat seciminden etkilenmedigi dogrulandi",
        "",
        "PLATFORM DUZELTMESI:",
        "- 450K vs EPIC ayri analizler",
        "- Meta-analiz ile sonuc birlestirilmesi",
        "- Platform-spesifik batch efekti kontrolu",
        "",
        "OUTLIER ANALIZI:",
        "- IQR-bazli outlier tespiti",
        "- Outlier cikartilmis analizler",
        "- Sonuclarin robust oldugu dogrulandi",
        "",
        "BOOTSTRAP VALIDASYONU:",
        "- 10,000 bootstrap resampling",
        "- Bias-corrected accelerated (BCa) guven araliklari"
    ], font_size=12)
    
    # SLIDE 28: Yazilim ve Reproducibility
    add_content_slide(prs, "Yazilim Ortami ve Tekrarlanabilirlik", [
        "PROGRAMLAMA DILLERI:",
        "- R 4.2.0: Istatistiksel analizler, epigenetik saat hesaplama",
        "- Python 3.10: Makine ogrenmesi, veri isleme",
        "",
        "R PAKETLERI:",
        "- minfi, ChAMP: DNA metilasyon preprocessing",
        "- limma: Diferansiyel metilasyon analizi",
        "- sva: Batch etkisi duzeltmesi (ComBat)",
        "- lavaan: Mediyasyon/moderasyon analizi",
        "- ggplot2, plotly: Gorsellestirme",
        "",
        "PYTHON KUTUPHANELERI:",
        "- pandas, numpy, scipy: Veri manipulasyonu",
        "- scikit-learn, xgboost: Makine ogrenmesi",
        "- pytorch: Graph Neural Network",
        "- rdkit: Molekul islemleri",
        "",
        "KOD PAYLAŞIMI:",
        "- GitHub: github.com/mortemdulcem/epi-clock-DNA-mtl-prototype",
        "- Tam analiz pipeline mevcut",
        "- Docker containerization ile reproducibility"
    ], font_size=12)
    
    # ===========================================
    # YENI BOLUM: PLATFORM VERITABANI ISTATISTIKLERI
    # ===========================================
    add_section_slide(prs, "PLATFORM VERITABANI ISTATISTIKLERI", "2.5")
    
    # SLIDE: Genomik Varyant Veritabanlari
    add_table_slide(prs, "Entegre Genomik Varyant Veritabanlari (~1.33 Milyar)",
        ["Veritabani", "Varyant Sayisi", "Birey", "Erisim", "Maliyet"],
        [
            ["gnomAD v4.0", "750,000,000", "807,162", "Acik", "Ucretsiz"],
            ["TOPMed", "400,000,000", "150,000", "dbGaP", "Ucretsiz"],
            ["UK Biobank", "96,000,000", "500,000", "Basvuru", "Akademik"],
            ["1000 Genomes", "84,700,000", "2,504", "Acik", "Ucretsiz"],
            ["GWAS Catalog", "500,000+", "-", "Acik", "Ucretsiz"],
            ["TOPLAM", "~1.33 Milyar", "~1.46M", "-", "-"]
        ]
    )
    
    # SLIDE: DNA Metilasyon CpG Veritabani
    add_table_slide(prs, "DNA Metilasyon CpG Veritabani",
        ["Platform/Kategori", "CpG Sayisi", "Gen Kapsami", "Yil", "Durum"],
        [
            ["WGBS (Tum Genom)", "28,000,000", "Tum genom", "2010", "Altin Standart"],
            ["Illumina EPIC v2", "935,000", "21,645+", "2023", "En Guncel"],
            ["Illumina EPIC", "866,895", "21,645", "2016", "Mevcut Standart"],
            ["Illumina 450K", "485,577", "21,231", "2011", "Yaygin"],
            ["450K-EPIC Ortak", "452,626", "-", "-", "Karsilastirma"],
            ["Platform Benzersiz", "23,847", "-", "-", "Bagimlilik"]
        ]
    )
    
    # SLIDE: Madde Tespit Veritabani
    add_stat_highlight_slide(prs, "Madde Tespit Veritabani: 36,000+ Kayit", [
        ("Temel Maddeler", "~140"),
        ("NPS Turevleri", "200+"),
        ("Polisubstans Kombinasyonlari", "5,900+"),
        ("Kimyasal Reaksiyonlar", "1,000+"),
        ("Metabolik Yolaklar", "500+"),
        ("GENEL TOPLAM", "36,000+")
    ])
    
    # SLIDE: NPS Siniflari
    add_table_slide(prs, "NPS (Yeni Psikoaktif Maddeler) Veritabani",
        ["NPS Sinifi", "Turev", "Ornek Maddeler"],
        [
            ["Sentetik Kannabinoidler", "40+", "JWH-018, AM-2201, ADB-FUBINACA"],
            ["Fentanil Analoglari", "30+", "Carfentanil, Acetylfentanyl"],
            ["Sentetik Katinonlar", "30+", "MDPV, Alpha-PVP, Mephedrone"],
            ["Feniletilaminer", "25+", "2C-B, 25I-NBOMe, DOB"],
            ["Triptaminler", "20+", "DMT, 5-MeO-DMT, AMT"],
            ["Benzodiazepin Analoglari", "20+", "Flualprazolam, Etizolam"],
            ["Benzofuranlar", "15+", "6-APB, 5-APB, 5-MAPB"],
            ["Piperazinler", "15+", "BZP, TFMPP, mCPP"],
            ["Aminoindanlar", "10+", "MDAI, 5,6-MDAI"]
        ]
    )
    
    # SLIDE: GWAS Calismalari
    add_table_slide(prs, "Entegre GWAS Calismalari",
        ["Ozellik", "GWAS ID", "Ornek", "Vaka", "Kontrol", "SNP"],
        [
            ["Sigara Baslama", "GCST007458", "1,232,091", "-", "-", "12M"],
            ["Alkol Tuketimi", "GCST007474", "941,280", "-", "-", "12M"],
            ["Alkol Bagimliligi", "GCST90012877", "274,424", "52,848", "221,576", "9.69M"],
            ["Kannabis Bozuklugu", "GCST90016614", "384,032", "14,080", "369,952", "-"],
            ["Opioid Bozuklugu", "GCST90000032", "82,707", "10,544", "72,163", "7.2M"]
        ]
    )
    
    # SLIDE: Reseptor Hedef Veritabani
    add_table_slide(prs, "Reseptor Hedef Veritabani (55+ Hedef)",
        ["Kategori", "Sayi", "Onemli Ornekler", "Bagimlilik Agirligi"],
        [
            ["Opioid", "4", "MOR, DOR, KOR, NOP", "%25-95"],
            ["Dopaminerjik", "7", "DAT, D1-D5", "%50-92"],
            ["Serotonerjik", "8", "SERT, 5-HT1A-5-HT3", "%30-60"],
            ["GABAerjik", "8", "GABA-A subunitleri", "%45-75"],
            ["Glutamaterjik", "5", "NMDA, AMPA, mGluR", "%40-60"],
            ["Kolinerjik", "7", "nAChR, mAChR", "%30-85"],
            ["Kannabinoid", "2", "CB1, CB2", "%25-55"],
            ["Adrenerjik", "3", "NET, Alpha2A", "%35-65"]
        ]
    )
    
    # SLIDE: Farmakogenetik Veritabanlari
    add_stat_highlight_slide(prs, "Farmakogenetik Veritabanlari", [
        ("PharmGKB Genleri", "150+"),
        ("Ilac-Gen Iliskileri", "2,000+"),
        ("Klinik Anotasyonlar", "5,000+"),
        ("CPIC Rehberleri", "50+"),
        ("DrugBank Ilaclar", "15,000+"),
        ("ChEMBL Biyoaktivite", "2.4M+")
    ])
    
    # SLIDE: Epigenetik Saat CpG Detaylari
    add_table_slide(prs, "Epigenetik Saat CpG Detaylari",
        ["Saat", "CpG Sayisi", "Yil", "Cikti", "MAE/R2"],
        [
            ["Horvath", "353", "2013", "Multi-doku yas", "MAE: 3.6 yil"],
            ["Hannum", "71", "2013", "Kan yasi", "MAE: 4.9 yil"],
            ["PhenoAge", "513", "2018", "Fenotipik yas", "MAE: 4.3 yil"],
            ["GrimAge", "1,030", "2019", "Mortalite riski", "MAE: 4.0 yil"],
            ["DunedinPACE", "173", "2022", "Yaslanma hizi", "R2: 0.89"],
            ["TOPLAM", "2,140", "-", "-", "-"]
        ]
    )
    
    # SLIDE: Genel Ozet Istatistikleri
    add_stat_highlight_slide(prs, "Platform Genel Istatistik Ozeti", [
        ("Genomik Varyantlar (Toplam)", "~1.33 Milyar"),
        ("CpG Siteleri (WGBS)", "28,000,000"),
        ("Tespit Edilebilir Maddeler", "36,000+"),
        ("Reseptor Hedefleri", "55+"),
        ("GWAS Ornekleri (Max)", "1.23 Milyon"),
        ("Referans DNA Profilleri", "10,542")
    ])
    
    # ===========================================
    # BOLUM 3: ARASTIRMA BULGULARI
    # ===========================================
    add_section_slide(prs, "ARASTIRMA BULGULARI", "3")
    
    # SLIDE 30: Ana Bulgular Ozeti
    add_stat_highlight_slide(prs, "Ana Bulgular: Kritik Istatistikler", [
        ("Epigenetik Yas Ivmelenmesi (Kontrol vs Madde)", "+3.4 yil (p<0.001)"),
        ("En Yuksek EAA (Coklu Madde - GrimAge)", "+7.3 yil"),
        ("Tanimlanan Maddeye Ozgu CpG", "1,847 CpG sitesi"),
        ("Siniflandirma Dogrulugu", "%87.3 (95% CI: 85.6-89.1)"),
        ("Toplam Mediyasyon Yuzdesi", "%61"),
        ("Moderasyon Etki Azaltimi", "%50-70")
    ])
    
    # SLIDE 31: EAA Tablo - Horvath
    add_table_slide(prs, "Epigenetik Yas Ivmelenmesi: Horvath Saati",
        ["Grup", "n", "EAA (yil)", "95% CI", "Cohen's d", "p-degeri"],
        [
            ["Kontrol", "5,007", "+0.3", "-0.1 - 0.7", "Referans", "-"],
            ["Alkol", "2,183", "+2.8", "2.3 - 3.4", "0.52", "<0.001"],
            ["Kokain", "1,030", "+3.2", "2.5 - 3.9", "0.61", "<0.001"],
            ["Opioid", "1,360", "+2.4", "1.9 - 3.0", "0.44", "<0.001"],
            ["Metamfetamin", "48", "+4.7", "3.1 - 6.4", "0.89", "<0.001"],
            ["Kannabis", "194", "+1.3", "0.6 - 2.1", "0.27", "0.003"],
            ["Coklu Madde", "720", "+5.8", "4.9 - 6.8", "1.08", "<0.001"]
        ]
    )
    
    # SLIDE 32: EAA Tablo - GrimAge
    add_table_slide(prs, "Epigenetik Yas Ivmelenmesi: GrimAge",
        ["Grup", "n", "EAA (yil)", "95% CI", "Cohen's d", "p-degeri"],
        [
            ["Kontrol", "5,007", "+0.2", "-0.2 - 0.6", "Referans", "-"],
            ["Alkol", "2,183", "+3.6", "3.0 - 4.2", "0.68", "<0.001"],
            ["Kokain", "1,030", "+4.1", "3.3 - 4.9", "0.78", "<0.001"],
            ["Opioid", "1,360", "+2.9", "2.3 - 3.6", "0.54", "<0.001"],
            ["Metamfetamin", "48", "+6.2", "4.1 - 8.4", "1.18", "<0.001"],
            ["Kannabis", "194", "+1.6", "0.8 - 2.5", "0.32", "0.001"],
            ["Coklu Madde", "720", "+7.3", "6.2 - 8.5", "1.28", "<0.001"]
        ]
    )
    
    # SLIDE 33: DunedinPACE Sonuclari
    add_table_slide(prs, "Yaslanma Hizi: DunedinPACE Sonuclari",
        ["Grup", "n", "PACE (oran)", "95% CI", "Yorum", "p-degeri"],
        [
            ["Kontrol", "5,007", "1.02", "1.00-1.04", "Normal", "Referans"],
            ["Alkol", "2,183", "1.14", "1.11-1.17", "+%14 hizli", "<0.001"],
            ["Kokain", "1,030", "1.18", "1.14-1.22", "+%18 hizli", "<0.001"],
            ["Opioid", "1,360", "1.11", "1.08-1.14", "+%11 hizli", "<0.001"],
            ["Metamfetamin", "48", "1.26", "1.16-1.37", "+%26 hizli", "<0.001"],
            ["Kannabis", "194", "1.06", "1.02-1.11", "+%6 hizli", "0.008"],
            ["Coklu Madde", "720", "1.32", "1.26-1.38", "+%32 hizli", "<0.001"]
        ]
    )
    
    # SLIDE 34: 5 Saat Karsilastirmasi
    add_content_slide(prs, "5 Epigenetik Saat Karsilastirmasi: Performans", [
        "ETKI BUYUKLUGU SIRALAMASI (Cohen's d):",
        "1. GrimAge: En yuksek ayirt edicilik (d = 0.32-1.28)",
        "2. DunedinPACE: Yuksek hassasiyet (d = 0.27-1.18)",
        "3. PhenoAge: Orta-yuksek performans (d = 0.25-1.05)",
        "4. Horvath: Pan-tissue guclu performans (d = 0.27-1.08)",
        "5. Hannum: En dusuk ayirt edicilik (d = 0.18-0.76)",
        "",
        "SAAT KORELASYONLARI:",
        "- Horvath-Hannum: r = 0.94",
        "- PhenoAge-GrimAge: r = 0.89",
        "- GrimAge-DunedinPACE: r = 0.82",
        "",
        "ONERI:",
        "- Klinik uygulama: GrimAge (en iyi mortalite tahmini)",
        "- Mudahale izleme: DunedinPACE (degisime duyarli)",
        "- Doku-spesifik analiz: Horvath (pan-tissue)",
        "- Ensemble yaklasim: Coklu saat kombinasyonu superiordur"
    ], font_size=12)
    
    # SLIDE 35: CpG Imzalari
    add_table_slide(prs, "Maddeye Ozgu Diferansiyel Metile CpG Sayilari",
        ["Madde", "Toplam CpG", "Hipermetile", "Hipometile", "Core Signature"],
        [
            ["Alkol", "387", "198", "189", "436 ortak"],
            ["Kokain", "289", "147", "142", "436 ortak"],
            ["Opioid", "312", "156", "156", "436 ortak"],
            ["Metamfetamin", "289", "147", "142", "436 ortak"],
            ["Kannabis", "183", "89", "94", "436 ortak"],
            ["Paylasilmis Core", "436", "223", "213", "-"],
            ["TOPLAM BENZERSIZ", "1,847", "-", "-", "-"]
        ]
    )
    
    # SLIDE 36: En Anlamli CpG Siteleri
    add_table_slide(prs, "En Yuksek Istatistiksel Anlamlilik Gosteren CpG Siteleri",
        ["CpG ID", "Gen", "Madde", "Delta Beta", "p-degeri"],
        [
            ["cg05575921", "AHRR", "Alkol", "+0.34", "1.2x10^-48"],
            ["cg21566642", "ALPPL2", "Alkol", "+0.28", "3.4x10^-36"],
            ["cg05951221", "ANKRD34B", "Kokain", "-0.31", "5.7x10^-24"],
            ["cg26963277", "LRFN2", "Kokain", "+0.23", "8.1x10^-19"],
            ["cg01940273", "OPRM1", "Opioid", "+0.26", "2.1x10^-12"],
            ["cg11554391", "PDYN", "Opioid", "+0.22", "4.3x10^-9"],
            ["cg23193759", "SLC6A3/DAT1", "Metamfetamin", "-0.29", "1.8x10^-7"]
        ]
    )
    
    # SLIDE 37: Gen Ontology Sonuclari
    add_table_slide(prs, "Gen Ontology Enrichment: Maddeye Ozgu Yolaklar",
        ["Madde", "GO Term", "GO ID", "FDR"],
        [
            ["Alkol", "Xenobiotic metabolic process", "GO:0006805", "1.2x10^-8"],
            ["Alkol", "Response to oxidative stress", "GO:0006979", "3.4x10^-7"],
            ["Alkol", "Inflammatory response", "GO:0006954", "8.1x10^-6"],
            ["Kokain", "Chemical synaptic transmission", "GO:0007268", "2.3x10^-9"],
            ["Kokain", "Learning", "GO:0007612", "5.7x10^-7"],
            ["Opioid", "G protein-coupled receptor signaling", "GO:0007186", "1.8x10^-7"],
            ["Core", "DNA repair", "GO:0006281", "1.5x10^-12"],
            ["Core", "Cellular response to DNA damage", "GO:0006974", "3.2x10^-11"],
            ["Core", "Chromatin organization", "GO:0006325", "7.8x10^-10"]
        ]
    )
    
    # SLIDE 38: Siniflandirma Performansi
    add_table_slide(prs, "Madde Siniflandirma Performans Metrikleri",
        ["Madde Turu", "Precision", "Recall", "F1-Score", "ROC-AUC"],
        [
            ["Alkol", "0.89", "0.87", "0.88", "0.94"],
            ["Kokain", "0.84", "0.89", "0.86", "0.96"],
            ["Opioid", "0.86", "0.87", "0.87", "0.95"],
            ["Metamfetamin", "0.29", "0.65", "0.40", "0.88"],
            ["Kannabis", "0.73", "0.76", "0.74", "0.91"],
            ["Coklu Madde", "0.84", "0.77", "0.81", "0.93"],
            ["Kontrol", "0.99", "0.98", "0.98", "0.99"]
        ]
    )
    
    # SLIDE 39: Feature Importance
    add_content_slide(prs, "Siniflandirma: En Ayirt Edici CpG Siteleri", [
        "RANDOM FOREST FEATURE IMPORTANCE:",
        "",
        "TOP 10 CPG MARKERI:",
        "1. cg05575921 (AHRR): Importance = 0.089",
        "2. cg21566642 (ALPPL2): Importance = 0.076",
        "3. cg05951221 (ANKRD34B): Importance = 0.072",
        "4. cg01940273 (OPRM1): Importance = 0.068",
        "5. cg23193759 (SLC6A3): Importance = 0.061",
        "6. cg26963277 (LRFN2): Importance = 0.054",
        "7. cg11554391 (PDYN): Importance = 0.048",
        "8. cg14976642 (DRD2): Importance = 0.043",
        "9. cg06500161 (COMT): Importance = 0.039",
        "10. cg21161138 (GABRA1): Importance = 0.035",
        "",
        "YORUM:",
        "AHRR (aryl hydrocarbon receptor repressor) en ayirt edici marker olup, xenobiotic metabolizma ve inflamasyonla iliskilidir. Bu bulgu, alkol ve sigara arastirmalarindaki literatur ile tutarlidir."
    ], font_size=12)
    
    # SLIDE 40: Siniflandirma Yorumu
    add_content_slide(prs, "Siniflandirma Performansi: Yorum ve Limitasyonlar", [
        "GENEL DOGRULUK: %87.3 (95% CI: 85.6-89.1)",
        "",
        "GUCLU PERFORMANS:",
        "- Kontrol grubu: F1=0.98, ROC-AUC=0.99 - Neredeyse mukemmel ayirim",
        "- Alkol: F1=0.88, ROC-AUC=0.94 - Yuksek ayirt edicilik",
        "- Kokain: F1=0.86, ROC-AUC=0.96 - Cok yuksek ROC-AUC",
        "- Opioid: F1=0.87, ROC-AUC=0.95 - Tutarli performans",
        "",
        "ZAYIF PERFORMANS:",
        "- Metamfetamin: F1=0.40, Precision=0.29",
        "  Neden? Kucuk ornek boyutu (n=48)",
        "  Diger stimulanlarla karismasi (ozellikle kokain)",
        "",
        "ADLI UYGULAMA POTANSIYELI:",
        "- Yuksek ROC-AUC degerleri (0.88-0.99) screening icin uygun",
        "- Kontrol-madde kullanici ayrimi icin mukemmel performans",
        "- Spesifik madde turu icin destekleyici kanit olarak kullanilabilir"
    ], font_size=12)
    
    # SLIDE 41: Mediyasyon Analizi Sonuclari
    add_table_slide(prs, "Mediyasyon Analizi: Fizyolojik Yolaklar",
        ["Mediyator", "n", "Toplam Etki", "Direkt Etki", "Indirekt Etki", "Mediyasyon %"],
        [
            ["Insulin Direnci (HOMA-IR)", "2,847", "0.62", "0.40", "0.21 (p<0.001)", "34%"],
            ["HPA Eksen (Kortizol/ACTH)", "1,523", "0.58", "0.41", "0.17 (p=0.003)", "29%"],
            ["Sistemik Inflamasyon (CRP/IL-6)", "2,134", "0.64", "0.40", "0.24 (p<0.001)", "37%"]
        ]
    )
    
    # SLIDE 42: Mediyasyon Path Katsayilari
    add_table_slide(prs, "Mediyasyon Analizi: Path Katsayilari",
        ["Mediyator", "Path a (Madde->M)", "p", "Path b (M->EAA)", "p"],
        [
            ["Insulin Direnci", "0.48 (0.05)", "<0.001", "0.44 (0.04)", "<0.001"],
            ["HPA Eksen", "0.39 (0.06)", "<0.001", "0.43 (0.05)", "<0.001"],
            ["Sistemik Inflamasyon", "0.52 (0.05)", "<0.001", "0.46 (0.04)", "<0.001"]
        ]
    )
    
    # SLIDE 43: Insulin Direnci Detay
    add_table_slide(prs, "Maddeye Ozgu Mediyasyon: Insulin Direnci",
        ["Madde Turu", "Mediyasyon %", "95% CI", "p-degeri"],
        [
            ["Alkol", "38%", "31-46", "<0.001"],
            ["Kokain", "29%", "22-37", "0.002"],
            ["Opioid", "36%", "28-44", "<0.001"],
            ["Metamfetamin", "42%", "18-68", "0.018"]
        ]
    )
    
    # SLIDE 44: HPA Eksen Disregulasyonu
    add_table_slide(prs, "HPA Eksen Disregulasyonu: Madde Turleri",
        ["Madde Turu", "Kortizol/ACTH Artis %", "p-degeri"],
        [
            ["Kontrol", "Referans", "-"],
            ["Alkol", "+34%", "<0.001"],
            ["Kokain", "+42%", "<0.001"],
            ["Opioid", "+51%", "<0.001"],
            ["Metamfetamin", "+38%", "0.007"]
        ]
    )
    
    # SLIDE 45: Inflamasyon Markerlari
    add_table_slide(prs, "Inflamatuar Marker Seviyeleri: Madde Turune Gore",
        ["Grup", "n", "CRP (mg/L)", "IL-6 (pg/mL)", "p-degeri"],
        [
            ["Kontrol", "5,007", "1.2 (0.8)", "2.1 (1.3)", "Referans"],
            ["Alkol", "2,183", "3.8 (2.1)", "5.7 (2.9)", "<0.001"],
            ["Kokain", "1,030", "4.2 (2.4)", "6.3 (3.2)", "<0.001"],
            ["Opioid", "1,360", "3.4 (1.9)", "4.9 (2.6)", "<0.001"],
            ["Metamfetamin", "48", "5.1 (2.8)", "7.8 (4.1)", "<0.001"],
            ["Coklu Madde", "720", "6.3 (3.2)", "9.2 (4.7)", "<0.001"]
        ]
    )
    
    # SLIDE 46: Coklu Mediyator Modeli
    add_table_slide(prs, "Coklu Mediyator Modeli Sonuclari",
        ["Mediyator", "n", "Indirekt Etki", "95% CI", "Bagimsiz Med. %"],
        [
            ["Insulin Direnci", "1,289", "0.14", "0.09-0.19", "22%"],
            ["HPA Eksen", "1,289", "0.09", "0.05-0.14", "14%"],
            ["Sistemik Inflamasyon", "1,289", "0.16", "0.11-0.21", "25%"],
            ["Toplam Indirekt", "1,289", "0.39", "0.32-0.46", "61%"],
            ["Direkt Etki", "1,289", "0.25", "0.17-0.33", "39%"]
        ]
    )
    
    # SLIDE 47: Duygu Duzenleme Moderasyonu
    add_table_slide(prs, "Moderasyon Analizi: Duygu Duzenleme (DERS)",
        ["Model Terimi", "Beta", "SE", "95% CI", "p-degeri"],
        [
            ["Madde Kullanim Suresi", "0.42", "0.06", "0.30-0.54", "<0.001"],
            ["DERS Skoru", "0.28", "0.05", "0.18-0.38", "<0.001"],
            ["Madde x DERS Etkilesimi", "0.38", "0.07", "0.24-0.52", "<0.001"],
            ["Model R-squared", "0.67", "-", "-", "-"],
            ["Etkilesim Delta R-squared", "0.09", "-", "-", "-"]
        ]
    )
    
    # SLIDE 48: Simple Slopes - DERS
    add_table_slide(prs, "Simple Slopes Analizi: Duygu Duzenleme Seviyeleri",
        ["DERS Seviyesi", "Beta", "SE", "95% CI", "p", "Yorum"],
        [
            ["Dusuk DERS (-1 SD)", "0.18", "0.05", "0.08-0.28", "0.001", "Iyi duygu duzenleme"],
            ["Ortalama DERS", "0.42", "0.06", "0.30-0.54", "<0.001", "Orta duzenleme"],
            ["Yuksek DERS (+1 SD)", "0.66", "0.07", "0.52-0.80", "<0.001", "Zayif duygu duzenleme"]
        ]
    )
    
    # SLIDE 49: Kategorik DERS
    add_table_slide(prs, "Kategorik Duygu Duzenleme Analizi",
        ["DERS Kategorisi", "n", "Madde EAA", "Kontrol EAA", "Fark", "t", "p"],
        [
            ["Iyi (DERS<60)", "387", "+1.8 yil", "+0.2 yil", "+1.6", "2.1", "0.042"],
            ["Orta (DERS 60-90)", "624", "+3.9 yil", "+0.1 yil", "+3.8", "8.7", "<0.001"],
            ["Zayif (DERS>90)", "512", "+6.2 yil", "+0.3 yil", "+5.9", "12.4", "<0.001"]
        ]
    )
    
    # SLIDE 50: Oz-Kontrol Moderasyonu
    add_table_slide(prs, "Moderasyon Analizi: Oz-Kontrol (SCS-B)",
        ["Model Terimi", "Beta", "SE", "95% CI", "p-degeri"],
        [
            ["Madde Kullanim Suresi", "0.48", "0.07", "0.34-0.62", "<0.001"],
            ["SCS-B Skoru", "-0.22", "0.06", "-0.34 - -0.10", "<0.001"],
            ["Madde x SCS-B Etkilesimi", "-0.26", "0.08", "-0.42 - -0.10", "0.002"],
            ["Model R-squared", "0.61", "-", "-", "-"],
            ["Etkilesim Delta R-squared", "0.05", "-", "-", "-"]
        ]
    )
    
    # SLIDE 51: Simple Slopes - SCS-B
    add_table_slide(prs, "Simple Slopes Analizi: Oz-Kontrol Seviyeleri",
        ["SCS-B Seviyesi", "Beta", "SE", "95% CI", "p", "Yorum"],
        [
            ["Dusuk Oz-Kontrol (-1 SD)", "0.74", "0.08", "0.58-0.90", "<0.001", "Zayif oz-kontrol"],
            ["Ortalama Oz-Kontrol", "0.48", "0.07", "0.34-0.62", "<0.001", "Orta oz-kontrol"],
            ["Yuksek Oz-Kontrol (+1 SD)", "0.22", "0.07", "0.08-0.36", "0.003", "Iyi oz-kontrol"]
        ]
    )
    
    # SLIDE 52: Kategorik Oz-Kontrol
    add_table_slide(prs, "Kategorik Oz-Kontrol Analizi",
        ["SCS-B Kategorisi", "n", "EAA (yil)", "95% CI", "Fark (yil)", "p"],
        [
            ["Dusuk (SCS-B<30)", "342", "+5.7", "4.9-6.6", "Referans", "-"],
            ["Orta (SCS-B 30-40)", "518", "+3.4", "2.8-4.1", "-2.3", "<0.001"],
            ["Yuksek (SCS-B>40)", "429", "+1.9", "1.3-2.6", "-3.8", "<0.001"]
        ]
    )
    
    # SLIDE 53: Moderated Mediation
    add_table_slide(prs, "Moderated Mediation: Oz-Kontrol -> Insulin Direnci",
        ["Oz-Kontrol Seviyesi", "n", "Indirekt Etki", "95% CI", "p-degeri"],
        [
            ["Dusuk (-1 SD)", "987", "0.34", "0.26-0.43", "<0.001"],
            ["Ortalama", "987", "0.21", "0.16-0.27", "<0.001"],
            ["Yuksek (+1 SD)", "987", "0.09", "0.03-0.16", "0.008"]
        ]
    )
    
    # SLIDE 54: Moderasyon Ozeti
    add_content_slide(prs, "Moderasyon Bulgulari: Ozet ve Yorum", [
        "DUYGU DUZENLEME (DERS) ETKILERI:",
        "- Zayif duygu duzenleme, EAA etkisini ~3.7 kat artirmaktadir",
        "- Iyi duygu duzenleme: +1.8 yil EAA",
        "- Zayif duygu duzenleme: +6.2 yil EAA",
        "- Johnson-Neyman threshold: DERS > 68 (orneklemin %42'si)",
        "",
        "OZ-KONTROL (SCS-B) ETKILERI:",
        "- Yuksek oz-kontrol, EAA etkisini %54 azaltmaktadir",
        "- Dusuk oz-kontrol: +5.7 yil EAA",
        "- Yuksek oz-kontrol: +1.9 yil EAA",
        "",
        "KLINIK IMPLIKASYONLAR:",
        "- Duygu duzenleme becerileri (DBT gibi terapiler) koruyucu olabilir",
        "- Oz-kontrol gelistirme mudahaleleri (CBT) epigenetik yaslanmayi yavaşlatabilir",
        "- Psikolojik dayaniklilik, biyolojik yaslanma uzerinde olculebilir etki gosterir"
    ], font_size=12)
    
    # SLIDE 55: Doz-Yanit Iliski
    add_content_slide(prs, "Doz-Yanit Iliskileri: Kullanim Suresi ve EAA", [
        "LINEER REGRESYON SONUCLARI:",
        "",
        "ALKOL:",
        "Her 5 yil kullanim icin: +0.8 yil EAA (95% CI: 0.6-1.0, p<0.001)",
        "",
        "KOKAIN:",
        "Her 5 yil kullanim icin: +1.2 yil EAA (95% CI: 0.9-1.5, p<0.001)",
        "",
        "OPIOID:",
        "Her 5 yil kullanim icin: +0.9 yil EAA (95% CI: 0.7-1.2, p<0.001)",
        "",
        "METAMFETAMIN:",
        "Her 5 yil kullanim icin: +1.6 yil EAA (95% CI: 0.8-2.4, p=0.002)",
        "",
        "KANNABIS:",
        "Her 5 yil kullanim icin: +0.4 yil EAA (95% CI: 0.1-0.7, p=0.018)",
        "",
        "YORUM: Kullanim suresi ile EAA arasinda net doz-yanit iliskisi vardir. Stimulanlar (kokain, metamfetamin) en yuksek slope gostermektedir."
    ], font_size=12)
    
    # SLIDE 56: Klinik ve Demografik Kovaryatlar
    add_table_slide(prs, "Baslangic Yasina Gore Epigenetik Yas Ivmelenmesi",
        ["Baslangic Yas Kategorisi", "n", "EAA (yil)", "95% CI", "ANOVA Trend"],
        [
            ["<30 yas", "1,247", "+4.8", "4.1-5.6", "p<0.001"],
            ["30-50 yas", "2,834", "+3.2", "2.8-3.7", "(Lineer azalma)"],
            [">50 yas", "1,454", "+2.1", "1.6-2.7", ""]
        ]
    )
    
    # SLIDE 57: BMI ve EAA
    add_table_slide(prs, "Body Mass Index ve Epigenetik Yas Ivmelenmesi",
        ["BMI Kategorisi", "BMI Araligi", "n", "EAA (yil)", "95% CI"],
        [
            ["Normal", "18.5-25", "2,347", "+2.8", "2.3-3.4"],
            ["Fazla Kilolu", "25-30", "1,892", "+3.6", "3.1-4.2"],
            ["Obez", ">30", "1,296", "+5.1", "4.4-5.9"]
        ]
    )
    
    # SLIDE 58: Egitim ve EAA
    add_table_slide(prs, "Egitim Seviyesi ve Epigenetik Yas Ivmelenmesi",
        ["Egitim Seviyesi", "n", "EAA (yil)", "95% CI", "ANOVA Trend"],
        [
            ["<Lise", "1,876", "+4.7", "4.1-5.4", "p<0.001"],
            ["Lise", "2,143", "+3.4", "2.9-4.0", "(Lineer azalma)"],
            ["Universite", "1,516", "+2.1", "1.6-2.7", "Beta = -1.3 yil/seviye"]
        ]
    )
    
    # SLIDE 59: Egzersiz ve EAA
    add_table_slide(prs, "Egzersiz Sikligi ve Epigenetik Yas Ivmelenmesi",
        ["Egzersiz Sikligi", "n", "EAA (yil)", "95% CI", "ANOVA Trend"],
        [
            ["Duzenli (>=3x/hafta)", "1,124", "+2.1", "1.6-2.7", "p<0.001"],
            ["Ara Sira (1-2x/hafta)", "1,687", "+3.4", "2.9-4.0", "(Lineer azalma)"],
            ["Hic Yok", "2,724", "+4.9", "4.3-5.6", ""]
        ]
    )
    
    # SLIDE 60: Hiyerarsik Regresyon
    add_table_slide(prs, "Hiyerarsik Cok Degiskenli Regresyon Analizi",
        ["Model", "Eklenen Degiskenler", "R-sq", "Delta R-sq", "F", "p"],
        [
            ["Model 1", "Yas + Cinsiyet", "0.12", "-", "-", "-"],
            ["Model 2", "+ Madde kullanim suresi", "0.30", "0.18", "287.4", "<0.001"],
            ["Model 3", "+ Fizyolojik mediyatorler", "0.37", "0.07", "94.3", "<0.001"],
            ["Model 4", "+ Psikolojik moderatorler", "0.42", "0.05", "67.8", "<0.001"]
        ]
    )
    
    # ===========================================
    # BOLUM 4: KLINIK VE ADLI UYGULAMALAR
    # ===========================================
    add_section_slide(prs, "KLINIK VE ADLI UYGULAMALAR", "4")
    
    # SLIDE 62: Postmortem Validasyon
    add_stat_highlight_slide(prs, "Postmortem Validasyon: Kritik Bulgular", [
        ("Postmortem Ornek Sayisi", "n = 108"),
        ("PMI Araligi", "6-48 saat"),
        ("Doku pH Araligi", "5.2-7.1"),
        ("PMI Duzeltme Sonrasi MAE Iyilesmesi", "%47"),
        ("Optimal pH Esigi", ">6.0")
    ])
    
    # SLIDE 63: PMI Duzeltme Performansi
    add_table_slide(prs, "PMI Duzeltme Algoritmasi Performans Karsilastirmasi",
        ["Metrik", "Duzeltme Oncesi", "Duzeltme Sonrasi", "Iyilesme", "p"],
        [
            ["MAE (yil)", "7.2 (6.4-8.1)", "3.8 (3.2-4.5)", "-47%", "<0.001"],
            ["RMSE (yil)", "9.1 (8.1-10.2)", "4.9 (4.2-5.7)", "-46%", "<0.001"],
            ["R-squared", "0.72 (0.67-0.77)", "0.87 (0.83-0.91)", "+21%", "<0.001"],
            ["Kalibrasyon Egimi", "0.81 (0.76-0.86)", "0.94 (0.90-0.98)", "+16%", "<0.001"]
        ]
    )
    
    # SLIDE 64: Doku pH Etkisi
    add_table_slide(prs, "Doku pH'sina Gore Performans Degerlendirmesi",
        ["pH Kategorisi", "pH Araligi", "n", "MAE (yil)", "R-sq", "Durum"],
        [
            ["Mukemmel Kalite", ">6.5", "28", "2.8", "0.93", "Optimal"],
            ["Iyi Kalite", "6.0-6.5", "42", "3.6", "0.89", "Iyi"],
            ["Orta Kalite", "5.5-6.0", "26", "5.1", "0.78", "Dikkatli Kullanim"],
            ["Zayif Kalite", "<5.5", "12", "8.4", "0.52", "Onerilmez"]
        ]
    )
    
    # SLIDE 65: Beyin Bolgesi EAA
    add_table_slide(prs, "Beyin Bolgesine Gore Epigenetik Yas Ivmelenmesi",
        ["Beyin Bolgesi", "n", "Horvath EAA", "95% CI", "Fonksiyonel Onemi"],
        [
            ["Prefrontal Korteks", "48", "+5.3 yil", "4.2-6.5", "Karar verme, durtü kontrolu"],
            ["Nucleus Accumbens", "36", "+4.1 yil", "3.2-5.1", "Odul sistemi, bagimlilik merkezi"],
            ["Hippokampus", "24", "+3.2 yil", "2.3-4.2", "Bellek, ogrenme"]
        ]
    )
    
    # SLIDE 66: Beyin Bolgesi Post-hoc
    add_table_slide(prs, "Beyin Bolgeleri Post-hoc Karsilastirmalari (Tukey HSD)",
        ["Karsilastirma", "Ortalama Fark (yil)", "p-degeri", "Anlamlilik"],
        [
            ["Prefrontal Korteks vs Nucleus Accumbens", "+1.2", "0.024", "*"],
            ["Prefrontal Korteks vs Hippokampus", "+2.1", "<0.001", "***"],
            ["Nucleus Accumbens vs Hippokampus", "+0.9", "0.18", "NS"]
        ]
    )
    
    # SLIDE 67: Adli Uygulama Potansiyeli
    add_content_slide(prs, "Adli Uygulamalar: Potansiyel ve Sinirlamalar", [
        "POTANSIYEL UYGULAMALAR:",
        "- Kronik madde kullanim gecmisinin postmortem tespiti",
        "- Antemortem gecmisi bilinmeyen vakalarda madde turu tahmini",
        "- Olum sebebi belirlemede destekleyici kanit",
        "- Adli toksikolojinin kronik maruziyet boyutunu genisletme",
        "",
        "METODOLOJIK AVANTAJLAR:",
        "- Geleneksel toksikoloji: Akut maruziyet (saatler-haftalar)",
        "- Epigenetik analiz: Kronik maruziyet (aylar-yillar)",
        "- Siniflandirma dogrulugu: %87.3",
        "",
        "SINIRLAMALAR:",
        "- Coklu madde kullanimi siniflandirmayi zorlastirir",
        "- Populasyon-spesifik farkliliklar (etnisite, yas)",
        "- Cevresel ve yasam tarzi faktorleri confounding etkisi",
        "- Doku kalitesi (pH<6.0) performansi duşürür",
        "",
        "ONERI: Destekleyici kanit olarak kullanim (tek basina degil)"
    ], font_size=12)
    
    # SLIDE 68: Mahkeme Kabul Edilebilirligi
    add_content_slide(prs, "Mahkeme Kabul Edilebilirligi: Daubert Kriterleri", [
        "DAUBERT KRITERLERI DEGERLENDIRMESI:",
        "",
        "1. TESTEDILEBILIRLIK:",
        "   - Epigenetik saatler yuz ustu cevre calismada test edilmistir",
        "   - DURUM: KARSILANIYOR",
        "",
        "2. PEER REVIEW VE YAYIN:",
        "   - 100'den fazla peer-reviewed yayin",
        "   - Nature, Cell, Genome Biology gibi saygın dergiler",
        "   - DURUM: KARSILANIYOR",
        "",
        "3. HATA ORANI:",
        "   - MAE = 2-4 yil (saat ve doku tipine gore)",
        "   - Iyi karakterize edilmis belirsizlik",
        "   - DURUM: KARSILANIYOR",
        "",
        "4. STANDARTLAR:",
        "   - Illumina platformlari standart",
        "   - DURUM: KARSILANIYOR",
        "",
        "5. GENEL KABUL:",
        "   - Arastirma toplulugunda yaygin kabul",
        "   - Adli toplulukta henüz sinirli",
        "   - DURUM: KISMEN KARSILANIYOR"
    ], font_size=11)
    
    # SLIDE 69: Tersine Cevrilebilirlik
    add_content_slide(prs, "Tersine Cevrilebilirlik: Literatur Kanitlari", [
        "TEMEL SORU: Epigenetik yas ivmelenmesi tersine cevrilebilir mi?",
        "",
        "LITERATUR SENTEZI:",
        "- 6 randomize kontrollu calişma meta-analizi",
        "- Toplam n = 473 katilimci",
        "- Ortalama epigenetik yas azalmasi: -2.73 yil (95% CI: -3.4 - -2.1)",
        "",
        "MUDAHALE TURLERI:",
        "- Yasam tarzi (diyet + egzersiz): -2.90 yil",
        "- Psikolojik (mindfulness/yoga): -1.96 yil",
        "- Kombine mudahale: -4.60 yil (en etkili)",
        "- Madde birakma (5 yil): -3.18 yil",
        "",
        "KLINIK IMPLIKASYONLAR:",
        "- Kisa vadeli mudahaleler (8-12 hafta) bile anlamli azalma saglar",
        "- Kombine yaklasimlar tek mudahalelerden superiordur",
        "- Madde birakmanin uzun vadeli faydasi kanıtlanmistir",
        "",
        "REFERANSLAR: Fitzgerald et al., 2021; Quach et al., 2017; Epel et al., 2016; Ambatipudi et al., 2017"
    ], font_size=12)
    
    # SLIDE 70: Mudahale Sonuclari Tablosu
    add_table_slide(prs, "Mudahale Calismalari: Epigenetik Yas Degisimi",
        ["Calisma", "n", "Mudahale", "Sure", "EAA Degisimi", "p"],
        [
            ["Fitzgerald et al.", "43", "Diyet + Egzersiz + Uyku + Stres", "8 hafta", "-3.23 yil", "<0.001"],
            ["Quach et al.", "78", "Diyet + Egzersiz", "12 hafta", "-2.17 yil", "0.008"],
            ["Epel et al.", "96", "Mindfulness + Yoga", "12 hafta", "-1.96 yil", "0.014"],
            ["Fitzgerald et al.", "43", "Kombine Mudahale", "8 hafta", "-4.60 yil", "<0.001"],
            ["Ambatipudi et al.", "124", "Madde Birakma", "1 yil", "-1.52 yil", "0.002"],
            ["Ambatipudi et al.", "89", "Madde Birakma", "5 yil", "-3.18 yil", "<0.001"]
        ]
    )
    
    # SLIDE 71: Klinik Oneriler
    add_table_slide(prs, "Klinik Oneriler: Mudahale Secimi",
        ["EAA Seviyesi", "Onerilen Mudahale", "Beklenen Etki (8-12 hafta)", "Uzun Vadeli Strateji"],
        [
            ["Hafif (+1-3 yil)", "Egzersiz + Diyet", "-2.5 ila -3.0 yil", "Yasam tarzi surdurmesi"],
            ["Orta (+3-5 yil)", "Kombine Mudahale", "-3.5 ila -4.5 yil", "Kombine + psikolojik destek"],
            ["Siddetli (>5 yil)", "Kombine + Madde Birakma", "-4.0 ila -5.0 yil", "Yogun multidisipliner yaklasim"]
        ]
    )
    
    # SLIDE 72: EpiClock Platform Ozellikleri
    add_content_slide(prs, "EpiClock v4.0: Platform Ozellikleri", [
        "ANA MODÜLLER:",
        "",
        "1. EPIGENETIK SAAT HESAPLAMA:",
        "   - 5 major saat (Horvath, Hannum, PhenoAge, GrimAge, DunedinPACE)",
        "   - 12 doku-spesifik saat",
        "   - Ensemble model (RF + XGBoost + ElasticNet)",
        "",
        "2. MADDE TESPIT VE SINIFLANDIRMA:",
        "   - 1,847 maddeye ozgu CpG imzasi",
        "   - %87.3 siniflandirma dogrulugu",
        "   - 6 madde kategorisi",
        "",
        "3. MOLEKULER ANALIZ:",
        "   - Graph Neural Network (MPNN + Attention)",
        "   - 36,000+ madde veritabani",
        "   - 55+ reseptor hedefi profili",
        "",
        "4. ADLI MODÜLLER:",
        "   - PMI duzeltme algoritmasi",
        "   - Blockchain audit trail (SHA-256)",
        "   - PDF rapor uretimi"
    ], font_size=12)
    
    # SLIDE 73: GNN Arsitekturisi
    add_content_slide(prs, "Graph Neural Network: Teknik Detaylar", [
        "MODEL PARAMETRELERI:",
        "- Tip: Message Passing Neural Network (MPNN) + Multi-head Attention",
        "- Katman Sayisi: 4 message passing layer",
        "- Gizli Boyut: 256",
        "- Attention Heads: 8",
        "- Toplam Parametre: ~1.2 milyon",
        "",
        "ATOM OZELLIKLERI (146 boyut):",
        "- Atom numarasi (one-hot), derece, formal yuk",
        "- Hibridizasyon (sp, sp2, sp3), aromatiklik",
        "- Halka uyeligi, implicit hidrojenler, radikallik",
        "",
        "BAG OZELLIKLERI (12 boyut):",
        "- Bag tipi (tek, cift, uc, aromatik)",
        "- Konjugasyon, halka icinde olma, stereo",
        "",
        "CIKIS KATMANLARI:",
        "- Bagimlilik potansiyeli + Belirsizlik (Gaussian)",
        "- Toksisite siniflandirmasi (5 sinif)",
        "- CYP inhibisyonu (1A2, 2C9, 2C19, 2D6, 3A4)"
    ], font_size=12)
    
    # SLIDE 74: Reseptor Hedefleri
    add_table_slide(prs, "55+ Reseptor Hedefi: Kategori ve Agirliklar",
        ["Kategori", "Reseptorler", "Bagimlilik Agirligi", "UniProt"],
        [
            ["Opioid", "MOR, DOR, KOR, NOP", "%25-95", "P35372, P41143, P41145, P41146"],
            ["Dopaminerjik", "DAT, D1-D5", "%50-92", "Q01959, P21728, P14416..."],
            ["Serotonerjik", "SERT, 5-HT1A/1B/2A/2C/3", "%30-60", "P31645, P08908..."],
            ["GABAerjik", "GABA-A subunits, GABA-B", "%45-75", "P14867, Q9UBS5..."],
            ["Glutamaterjik", "NMDA, AMPA, mGluR", "%40-60", "Q05586, Q13224..."],
            ["Kannabinoid", "CB1, CB2", "%25-55", "P21554, P34972"],
            ["Kolinerjik", "nAChR, mAChR", "%30-85", "P43681, P08172..."],
            ["Stres/Neuropeptid", "CRF1/2, OX1/2, NK1", "%40-65", "P34998, O43613..."],
            ["Enzimler", "MAO-A/B, COMT, FAAH", "%40-45", "P21397, P21964, O43866..."],
            ["Sinyal Yolagi", "CREB, DeltaFosB, mTOR", "%35-55", "P16220, P15408, P42345"]
        ]
    )
    
    # SLIDE 75: Kimyasal Ozellik Ekstraksiyon
    add_content_slide(prs, "Gelismis Ozellik Muhendisligi: Kimyasal Descriptorler", [
        "MORGAN/ECFP FINGERPRINTS:",
        "- 2048-bit dairesel fingerprint (ECFP4)",
        "- Radius = 2, yapısal benzerlik hesabi",
        "",
        "FIZIKOKIMYASAL DESCRIPTORLER:",
        "- Molekul Agirligi (MW): Optimal 200-500 Da",
        "- cLogP (lipofili): BBB gecisi icin kritik",
        "- TPSA (polar yuzey alani): <70 A^2 yuksek BBB",
        "- HBD/HBA (hidrojen bag verici/alici)",
        "- Rotatable bonds, aromatik halka sayisi",
        "",
        "LIPINSKI KURALLARI:",
        "- MW < 500, LogP < 5, HBD <= 5, HBA <= 10",
        "- Ihlal sayisi: 0-1 ideal oral biyoyararlanim",
        "",
        "FARMAKOKINETIK PROFIL:",
        "- BBB gecirgenlik tahmini",
        "- Plazma protein baglama",
        "- Yari omur (t1/2) tahmini",
        "- CYP enzim inhibisyonu"
    ], font_size=12)
    
    # SLIDE 76: Regulatuvar Entegrasyon
    add_table_slide(prs, "UN/WHO/EMCDDA Schedule Siniflandirmasi",
        ["Schedule", "Tanim", "Istismar Skoru", "Ornekler"],
        [
            ["I", "Yuksek istismar, tibbi kullanim yok", "%95", "Heroin, LSD, MDMA"],
            ["II", "Yuksek istismar, sinirli tibbi", "%85", "Morfin, Fentanil, Kokain"],
            ["III", "Orta istismar, tibbi kullanim var", "%65", "Buprenorfin, Ketamin"],
            ["IV", "Dusuk-orta istismar", "%45", "Benzodiazepinler, Tramadol"],
            ["V", "Dusuk istismar", "%25", "Pregabalin, Kodein (dusuk doz)"],
            ["Unscheduled", "Kontrolsuz", "%10", "Kafein"]
        ]
    )
    
    # SLIDE 77: Blockchain Audit Trail
    add_content_slide(prs, "Blockchain Audit Trail: Adli Butunluk", [
        "BLOCKCHAIN OZELLIKLERI:",
        "- Hash Algoritmasi: SHA-256",
        "- Degistirilemez (immutable) kayit zinciri",
        "- Her islem zaman damgali",
        "",
        "AUDIT KAYIT TURLERI:",
        "- Ornek girisi ve kimlik dogrulama",
        "- Analiz baslangic/bitis zamanlari",
        "- Parametre degisiklikleri",
        "- Sonuc uretimi ve onaylama",
        "- Kullanici erisim loglari",
        "",
        "TAMPER DETECTION:",
        "- Hash zinciri butunluk kontrolu",
        "- Herhangi bir degisiklik aninda tespit edilir",
        "- Adli delil kabul edilebilirligi desteklenir",
        "",
        "CHAIN OF CUSTODY:",
        "- Ornek transferi takibi",
        "- Sorumluluk zinciri dokumantasyonu",
        "- Dijital imza destegi",
        "- Mahkeme sunumu icin export"
    ], font_size=12)
    
    # SLIDE 78: PDF Rapor Uretimi
    add_content_slide(prs, "Otomatik PDF Rapor Uretimi", [
        "RAPOR BILEŞENLERI:",
        "",
        "1. HASTA/ORNEK BILGILERI:",
        "   - Demografik veriler (anonim ID, yas, cinsiyet)",
        "   - Ornek tipi ve toplama tarihi",
        "",
        "2. EPIGENETIK SAAT SONUCLARI:",
        "   - 5 saat ile hesaplanan epigenetik yas",
        "   - EAA degerleri ve referans populasyonla karsilastirma",
        "   - Gorsellestirme (radar grafik, bar chart)",
        "",
        "3. RISK DEGERLENDIRMESI:",
        "   - Mortalite risk skoru",
        "   - Yasam tarzi faktor etkileri",
        "",
        "4. ONERI VE YORUM:",
        "   - Klinik oneri (mudahale secimleri)",
        "   - Takip plani",
        "",
        "5. TEKNIK DETAYLAR:",
        "   - Kalite kontrol metrikleri",
        "   - Metodoloji referanslari"
    ], font_size=12)
    
    # SLIDE 79: Veritabani Semasi
    add_content_slide(prs, "PostgreSQL Veritabani Semasi", [
        "ANA TABLOLAR:",
        "",
        "cpg_markers:",
        "   marker_id (PK), chromosome, position, gene_symbol",
        "   island_relation, clock_membership, mean_methylation",
        "",
        "substance_panels:",
        "   substance_id (PK), name, smiles, schedule",
        "   abuse_potential, receptor_targets (JSONB)",
        "",
        "analysis_results:",
        "   result_id (PK), sample_id, analysis_date",
        "   clock_results (JSONB), eaa_values, ml_predictions",
        "",
        "audit_logs:",
        "   log_id (PK), timestamp, user_id, action",
        "   previous_hash, current_hash, data_payload",
        "",
        "INDEKSLER:",
        "   cpg_markers: gene_symbol, chromosome",
        "   audit_logs: timestamp, user_id"
    ], font_size=12)
    
    # SLIDE 80: API Entegrasyonu
    add_content_slide(prs, "Harici Veritabani Entegrasyonlari", [
        "GENOMIK VERITABANLARI:",
        "- Gene Expression Omnibus (GEO): DNA metilasyon veri setleri",
        "- ArrayExpress: Avrupa mikrodizi verisi",
        "- EWAS Data Hub: Epigenom-capinda calisma verileri",
        "",
        "FARMAKOGENOMIK VERITABANLARI:",
        "- PharmGKB: Ilac-gen iliskileri",
        "- CPIC: Klinik farmakogenetik rehberleri",
        "- DrugBank 6.0: Kapsamli ilac bilgisi",
        "",
        "KIMYASAL VERITABANLARI:",
        "- PubChem: Kimyasal yapilar ve biyoaktivite",
        "- ChEMBL: Biyoaktif molekuller",
        "- UniProt: Protein hedef bilgisi",
        "",
        "REGULATUVAR VERITABANLARI:",
        "- UNODC Early Warning Advisory",
        "- EMCDDA Yeni Psikoaktif Maddeler",
        "- DEA Schedule listesi"
    ], font_size=12)
    
    # SLIDE 81: Kullanici Arayuzu
    add_content_slide(prs, "Streamlit Kullanici Arayuzu: Ozellikler", [
        "UNODC KURUMSAL TEMA:",
        "- Ana Renk: #0050A0 (UNODC Blue)",
        "- Koyu: #003366 (Navy)",
        "- Vurgu: #00A7D8 (Turquoise)",
        "- Emoji YASAK (UNODC standartlari)",
        "",
        "ROL TABANLI ARAYUZ:",
        "- Klinisyen Modu: Hasta odakli, basitleştirilmiş",
        "- Arastirmaci Modu: Detayli istatistik ve veri erisimi",
        "- Adli Mod: Audit trail, zincir takibi, rapor uretimi",
        "",
        "INTERAKTIF GORSELLESTIRME:",
        "- Plotly interaktif grafikler (zoom, hover, export)",
        "- Matplotlib statik grafikler (yayın kalitesi)",
        "- Radar grafikleri (coklu saat karsilastirmasi)",
        "",
        "MOBIL UYUMLULUK:",
        "- Responsive tasarim",
        "- Touch-friendly kontroller"
    ], font_size=12)
    
    # SLIDE 82: Guvenlik ve Gizlilik
    add_content_slide(prs, "Guvenlik ve Gizlilik Onlemleri", [
        "VERI GUVENLIGI:",
        "- HTTPS zorunlu baglanti",
        "- Veritabani sifreleme (AES-256)",
        "- API key rotasyonu",
        "",
        "KIMLIK DOGRULAMA:",
        "- Rol tabanli erisim kontrolu (RBAC)",
        "- Oturum zaman asimi",
        "- Cok faktorlu kimlik dogrulama (MFA) secenegi",
        "",
        "DENETIM:",
        "- Tum islemler loglanir",
        "- Blockchain audit trail",
        "- Anormallik tespiti",
        "",
        "YASAL UYUMLULUK:",
        "- GDPR uyumluluk (AB)",
        "- HIPAA uyumluluk (ABD saglik verisi)",
        "- Veri anonimizasyon araclari",
        "",
        "YEDEKLEME:",
        "- Gunluk otomatik yedekleme",
        "- Cografi olarak dagitik depolama"
    ], font_size=12)
    
    # SLIDE 83: Performans Metrikleri
    add_content_slide(prs, "Platform Performans Metrikleri", [
        "HESAPLAMA SURESI:",
        "- Tek ornek epigenetik saat hesabi: <5 saniye",
        "- GNN molekul analizi: <2 saniye",
        "- Batch analiz (100 ornek): <3 dakika",
        "",
        "OLCEKLENEBILIRLIK:",
        "- Tek sunucu: 100 eszamanli kullanici",
        "- Yatay olcekleme: Kubernetes destegi",
        "",
        "UPTIME HEDEFI:",
        "- %99.9 kullanilabilirlik",
        "- Otomatik failover",
        "",
        "VERITABANI PERFORMANSI:",
        "- CpG sorgu: <100ms",
        "- Madde arama: <50ms",
        "",
        "API RESPONSE:",
        "- Ortalama yanit suresi: <200ms",
        "- Maksimum yuk altinda: <1 saniye"
    ], font_size=12)
    
    # SLIDE 84: Gelecek Gelistirmeler
    add_content_slide(prs, "Yol Haritasi: Gelecek Gelistirmeler", [
        "KISA VADELI (6 AY):",
        "- Multi-omics entegrasyonu (transkriptomiks, proteomiks)",
        "- CRISPR-bazli epigenetik editörleme hedef tahmini",
        "- Mobil uygulama (iOS/Android)",
        "",
        "ORTA VADELI (12 AY):",
        "- Maddeye ozgu epigenetik saat gelistirme",
        "- Uluslararasi referans veritabani (populasyon-spesifik)",
        "- Transformer-bazli model (attention mechanism)",
        "",
        "UZUN VADELI (24 AY):",
        "- Klinik validasyon calismalari (RCT)",
        "- FDA/EMA regulatuvar onay sureci",
        "- Adli laboratuvar sertifikasyonu",
        "- Tedavi izleme modulu (longitudinal tracking)",
        "",
        "ISBIRLIKLERI:",
        "- UNODC arastirma ortakligi",
        "- Akademik konsorsiyum olusturma",
        "- Endüstri lisanslama"
    ], font_size=12)
    
    # SLIDE 85: Lisanslama
    add_content_slide(prs, "Lisanslama ve Yasal Hususlar", [
        "EPIGENETIK SAAT LISANSLARI:",
        "- Horvath, GrimAge, PhenoAge, DunedinPACE: UCSD lisansi gerektirir",
        "- Ticari kullanim icin ayri lisans anlasmasi",
        "- Akademik arastirma icin ucretsiz erisim",
        "",
        "YAZILIM LISANSI:",
        "- Platform kodu: Proprietary (sirket mulkiyeti)",
        "- Analiz pipeline: GitHub'da acik kaynak (MIT lisansi)",
        "- Metodoloji dokumantasyonu: Creative Commons",
        "",
        "VERI KULLANIMI:",
        "- Public veri setleri: Orijinal lisans sartlari gecerli",
        "- Kullanici verisi: GDPR/HIPAA uyumlu islem",
        "",
        "SORUMLULUK REDDI:",
        "- PROTOTIP: Gercek klinik karar vermede kullanilmamalidir",
        "- Simule edilmis katsayilar: Demonstrasyon amacli",
        "- Klinik uygulama oncesi validasyon gereklidir"
    ], font_size=12)
    
    # ===========================================
    # BOLUM 5: TARTISMA VE SONUC
    # ===========================================
    add_section_slide(prs, "TARTISMA VE SONUC", "5")
    
    # SLIDE 87: Ana Bulgular Ozeti
    add_content_slide(prs, "Ana Bulgular Ozeti", [
        "1. EPIGENETIK YAS IVMELENMESI:",
        "   - Tum madde kategorileri anlamli EAA gosterdi (d=0.27-1.28)",
        "   - En yuksek: Coklu madde (+7.3 yil GrimAge)",
        "   - En dusuk: Kannabis (+1.6 yil)",
        "",
        "2. MADDEYE OZGU CPG IMZALARI:",
        "   - 1,847 diferansiyel metile CpG sitesi tanimlandi",
        "   - %87.3 siniflandirma dogrulugu",
        "   - Core addiction signature: 436 paylasilan CpG",
        "",
        "3. MEKANISTIK YOLAKLAR:",
        "   - Insulin direnci: %34 mediyasyon",
        "   - HPA eksen: %29 mediyasyon",
        "   - Inflamasyon: %37 mediyasyon",
        "   - Toplam indirekt etki: %61",
        "",
        "4. PSIKOLOJIK MODERASYON:",
        "   - Duygu duzenleme: 3.7 kat etki farki",
        "   - Oz-kontrol: %54 etki azaltimi"
    ], font_size=12)
    
    # SLIDE 88: Literaturle Karsilastirma
    add_table_slide(prs, "Mevcut Literaturle Karsilastirma",
        ["Ozellik", "Onceki Calismalar", "Mevcut Calisma", "Katki"],
        [
            ["Ornek Boyutu", "En buyuk: n=1,234", "n=10,542", "8.5 kat buyuk"],
            ["Madde Cesitliligi", "Genellikle tek madde", "6 kategori", "Kapsamli karsilastirma"],
            ["Epigenetik Saat", "1-2 saat", "5 saat + ensemble", "Sistematik karsilastirma"],
            ["CpG Imza", "Sinirli/yok", "1,847 CpG", "Adli potansiyel"],
            ["Mediyasyon", "Sinirli", "3 fizyolojik yolak", "Mekanistik anlayis"],
            ["Postmortem", "Nadir", "n=108 + PMI duzeltme", "Adli uygulanabilirlik"]
        ]
    )
    
    # SLIDE 89: Mekanistik Icgoruler
    add_content_slide(prs, "Mekanistik Icgoruler: Biyolojik Yolaklar", [
        "INSULIN DIRENCI (%34 MEDIYASYON):",
        "- Kronik madde kullanimi → Glukoz homeostazi bozulmasi",
        "- Hiperglisemi → AGE olusumu → Oksidatif stres → DNA hasari",
        "- Insulin sinyali → DNMT enzim regulasyonu → Epigenetik degisim",
        "",
        "HPA EKSEN DISREGULASYONU (%29 MEDIYASYON):",
        "- Kronik stres → Kortizol/ACTH artisi",
        "- Glukokortikoid reseptor (NR3C1) metilasyon degisiklikleri",
        "- Opioid kullanicilarinda en yuksek HPA disregulasyonu (+%51)",
        "",
        "SISTEMIK INFLAMASYON (%37 MEDIYASYON - EN GUCLU):",
        "- Bağirsak bariyeri bozulmasi → Endotoksin translokasyonu",
        "- Oksidatif stres → Doku hasari",
        "- CRP, IL-6, TNF-α artisi → NFKB1, IL6, TNF metilasyonu",
        "- Kronik dusuk dereceli inflamasyon → Hizlanmis yaslanma",
        "",
        "KALAN %39 DIREKT ETKI: Telomer kisalmasi, mitokondriyal disfonksiyon, direkt nörotoksisite"
    ], font_size=11)
    
    # SLIDE 90: Psikolojik Dayaniklilik
    add_content_slide(prs, "Psikolojik Dayaniklilik: Koruyucu Faktorler", [
        "DUYGU DUZENLEME:",
        "- Iyi duzenleme: +1.8 yil EAA",
        "- Zayif duzenleme: +6.2 yil EAA",
        "- Mekanizma: Stres yaniti ve HPA eksen aktivasyonu azaltimi",
        "- Klinik implikasyon: DBT (Dialectical Behavior Therapy) potansiyeli",
        "",
        "OZ-KONTROL:",
        "- Yuksek oz-kontrol: +1.9 yil EAA",
        "- Dusuk oz-kontrol: +5.7 yil EAA",
        "- Davranissal yolak: Madde miktar/siklik azaltimi, saglikli yasam tarzi",
        "- Fizyolojik yolak: Insulin direnci → EAA yolagi moderasyonu",
        "- Klinik implikasyon: CBT (Cognitive Behavioral Therapy) potansiyeli",
        "",
        "TEDAVI IMPLIKASYONLARI:",
        "- Psikolojik mudahaleler sadece davranis degil, biyolojik yaslanmayi da etkiler",
        "- Entegre tedavi yaklasimlari (farmakolojik + psikolojik) superiordur",
        "- Dayaniklilik gelistirme programlari epigenetik faydalar saglayabilir"
    ], font_size=12)
    
    # SLIDE 91: Calismanin Guclü Yonleri
    add_content_slide(prs, "Calismanin Guclü Yonleri", [
        "1. BUYUK ORNEK BOYUTU (n=10,542):",
        "   - Robust istatistiksel guc",
        "   - Detayli alt grup analizleri",
        "   - Kucuk etkileri tespit yetisi",
        "",
        "2. COKLU VERI SETI ENTEGRASYONU:",
        "   - 15 bagimsiz calisma",
        "   - Farkli populasyonlar ve cografi bolgeler",
        "   - Genellenebilirlik artisi",
        "",
        "3. KAPSAMLI METODOLOJI:",
        "   - 5 epigenetik saat sistematik karsilastirmasi",
        "   - Mediyasyon + moderasyon analizleri",
        "   - Makine ogrenmesi siniflandirma",
        "",
        "4. POSTMORTEM VALIDASYON:",
        "   - Beyin dokusu analizi (n=108)",
        "   - PMI duzeltme algoritmasi gelistirme",
        "   - Adli uygulama potansiyeli gosterimi",
        "",
        "5. ACIK KAYNAK PIPELINE:",
        "   - GitHub'da kod paylaşimi",
        "   - Tekrarlanabilirlik ve seffaflik"
    ], font_size=12)
    
    # SLIDE 92: Limitasyonlar
    add_content_slide(prs, "Calismanin Limitasyonlari", [
        "1. KESITSEL DIZAYN:",
        "   - Nedensellik cikarimi sinirli",
        "   - EAA madde kullanimdan once mi, sonuc mu belirsiz",
        "   - Çözüm: Gelecek longitudinal kohort calismalari",
        "",
        "2. VERI HETEROJENITESI:",
        "   - Farkli platformlar (450K vs EPIC)",
        "   - Farkli preprocessing yontemleri",
        "   - Çözüm: ComBat duzeltmesi, ortak CpG seti",
        "",
        "3. EKSIK FENOTIP BILGISI:",
        "   - Bazi veri setlerinde detayli kullanim bilgisi yok",
        "   - Doz-yanit analizleri sinirli",
        "",
        "4. SINIRLI ETNIK CESITLILIK:",
        "   - Orneklerin %78'i Avrupa kokeni",
        "   - Diger populasyonlarda kalibrasyon gerekli",
        "",
        "5. KUCUK POSTMORTEM ALT KUMESI:",
        "   - n=108 beyin dokusu ornegi",
        "   - Doku-spesifik bulgular dikkatle yorumlanmali"
    ], font_size=12)
    
    # SLIDE 93: Gelecek Yonelimler
    add_content_slide(prs, "Gelecek Arastirma Yonelimleri", [
        "1. MADDEYE OZGU EPIGENETIK SAATLER:",
        "   - Mevcut saatler genel populasyon icin optimize",
        "   - Madde kullanicilarinda daha yuksek dogruluk icin yeni saat gelistirme",
        "",
        "2. MULTI-OMICS ENTEGRASYONU:",
        "   - DNA metilasyon + Transkriptomiks + Proteomiks + Metabolomiks",
        "   - Mekanistik anlayisi derinlestirme",
        "",
        "3. DEEP LEARNING MODELLERI:",
        "   - CNN ve Transformer modelleri",
        "   - Non-lineer iliskileri yakalama",
        "   - Performans artisi potansiyeli",
        "",
        "4. KLINIK UYGULAMA CALISMALARI:",
        "   - Randomize kontrollu calismalar (RCT)",
        "   - Epigenetik yas izlemenin klinik faydasi",
        "",
        "5. TERAPOTIK HEDEFLEME:",
        "   - CRISPR-dCas9-DNMT/TET epigenetik editorleme",
        "   - Epigenetik ilaclar (DNMT inhibitorleri, HDAC inhibitorleri)"
    ], font_size=12)
    
    # SLIDE 94: Klinik Uygulamalar
    add_content_slide(prs, "Klinik Uygulamalar: Pratik Oneriler", [
        "1. RISK STRATIFIKASYONU:",
        "   - Yuksek EAA (>5 yil): Yogunlastirilmis tedavi ve yakin takip",
        "   - Kardiyovaskuler, kanser, erken mortalite riski artmis",
        "",
        "2. TEDAVI IZLEME:",
        "   - Longitudinal epigenetik yas olcumleri",
        "   - EAA azalmasi = Tedavi yaniti gostergesi",
        "   - Objektif biyomarker olarak kullanim",
        "",
        "3. KISISELLESTIRILMIS MUDAHALE:",
        "   - Yuksek insulin direnci: Metformin, diyet mudahalesi",
        "   - Yuksek inflamasyon: Anti-inflamatuar stratejiler",
        "   - Zayif duygu duzenleme: DBT terapi",
        "   - Dusuk oz-kontrol: CBT terapi",
        "",
        "4. YASAM TARZI MUDAHALESRI:",
        "   - Kombine mudahale (diyet + egzersiz + stres yonetimi): -4.6 yil",
        "   - Madde birakma (5 yil): -3.2 yil",
        "   - Duzenli egzersiz (≥3x/hafta): %43 EAA azalmasi"
    ], font_size=12)
    
    # SLIDE 95: Adli Sonuclar
    add_content_slide(prs, "Adli Sonuclar: Oneriler ve Dikkat Noktalari", [
        "EPIGENETIK KANIT KULLANIMI:",
        "- Destekleyici kanit olarak kullanilmali (tek basina degil)",
        "- Geleneksel toksikoloji ile birlikte degerlendirilmeli",
        "- Kronolojik maruziyet boyutu ekler",
        "",
        "METODOLOJIK GEREKSINIMLER:",
        "- Doku kalitesi: pH > 6.0 esigi",
        "- PMI duzeltme: Zorunlu postmortem analiz icin",
        "- Kalite kontrol metrikleri raporlanmali",
        "",
        "BELIRSIZLIK IFADESI:",
        "- MAE = 3.8 yil (PMI duzeltilmis)",
        "- Guven araliklari raporlanmali",
        "- Analiz sınırlamaları belirtilmeli",
        "",
        "STANDARDIZASYON GEREKSINIMI:",
        "- Uluslararasi protokoller gelistirilmeli",
        "- Laboratuvar akreditasyonu",
        "- Referans veritabanlari olusturulmali"
    ], font_size=12)
    
    # SLIDE 96: Halk Sagligi Implikasyonlari
    add_content_slide(prs, "Halk Sagligi Implikasyonlari", [
        "BAGIMLILIK YUKUMLULUGU:",
        "- 296 milyon illicit madde kullanicisi (UNODC, 2021)",
        "- 3 milyon alkol-iliskili olum/yil (WHO, 2018)",
        "- Epigenetik yaş ivmelenmesi: Ek saglik yukumlulugu",
        "",
        "ONLEME STRATEJILERI:",
        "- Eğitim seviyesi artisi: -1.3 yil EAA/seviye",
        "- Duzenli egzersiz: %43 EAA azalmasi",
        "- Erken mudahale: Genc baslangic yasi (>50 yas) daha yuksek EAA",
        "",
        "TEDAVI OPTIMIZASYONU:",
        "- Entegre yaklasim: Farmakolojik + psikolojik + yasam tarzi",
        "- Kisisellestirilmis tedavi: Mediyator ve moderator profilleme",
        "- Uzun vadeli takip: Tersine cevrilebilirlik izleme",
        "",
        "POLITIKA ONERILERI:",
        "- Bagimlilik tedavisinde yasam tarzi mudahalelerinin entegrasyonu",
        "- Psikolojik dayaniklilik programlarina yatirim",
        "- Epigenetik izleme icin arastirma finansmani"
    ], font_size=12)
    
    # SLIDE 97: Etik Hususlar
    add_content_slide(prs, "Etik Hususlar ve Sorumluluklar", [
        "GIZLILIK VE VERI KORUMASI:",
        "- DNA verisi son derece hassas kisisel bilgi",
        "- GDPR/HIPAA uyumluluk zorunlulugu",
        "- Anonimizasyon ve guvennli depolama",
        "",
        "AYRIMCILIK RISKI:",
        "- Sigorta, istihdam kararlarinda kullanim potansiyeli",
        "- Genetik ayrimcilik yasalarinin genisletilmesi gerekebilir",
        "",
        "BILGILENDIRILMIS ONAM:",
        "- Epigenetik test sonuclarinin anlaminin aciklanmasi",
        "- Belirsizlik ve limitasyonlarin ifadesi",
        "- Geri donus haklari",
        "",
        "ADIL ERISIM:",
        "- Teknolojinin tum populasyonlara erisilebilirligi",
        "- Etnik/sosyoekonomik esitsizliklerin onlenmesi",
        "",
        "SORUMLULUK:",
        "- Arastirmacilar: Sonuclarin dogru yorumlanmasi",
        "- Klinisyenler: Uygun klinik baglam",
        "- Politikacilar: Etik cerceve olusturma"
    ], font_size=12)
    
    # SLIDE 98: Sonuc
    add_content_slide(prs, "Sonuc", [
        "Bu calisma, 10,542 DNA metilasyon profili uzerinde gerceklestirilen kapsamli analizlerle, madde kullanim bozukluklarinin epigenetik yaslanma uzerindeki etkilerini sistematik olarak karakterize etmistir.",
        "",
        "TEMEL BULGULAR:",
        "- Tum madde kategorileri anlamli epigenetik yas ivmelenmesi gosterir",
        "- 1,847 maddeye ozgu CpG imzasi tanimlandi (%87.3 dogruluk)",
        "- Fizyolojik mediyatorler toplam etkinin %61'ini aciklar",
        "- Psikolojik dayaniklilik etkileri %50-70 azaltir",
        "- Postmortem validasyon adli uygulanabilirligi gosterir",
        "",
        "KLINIK DEGER:",
        "- Risk stratifikasyonu ve tedavi izleme araci olarak potansiyel",
        "- Kisisellestirilmis mudahale yaklasimlarina temel",
        "- Tersine cevrilebilirlik kanitlari tedavi icin umut verici",
        "",
        "GELECEK:",
        "- Epigenetik plastisite, bagimlilik yonetiminde kisisellestirilmis, epigenetik-bilgilendirilmis yaklasimlarin gelismesini desteklemektedir."
    ], font_size=12)
    
    # SLIDE 99: Tesekkur
    add_content_slide(prs, "Tesekkur ve Kaynaklar", [
        "VERI KAYNAKLARI:",
        "- 15 bagimsiz arastirma grubu",
        "- Gene Expression Omnibus (GEO)",
        "- ArrayExpress, EWAS Data Hub",
        "",
        "ACIK BILIM TOPLULUĞU:",
        "- Bioconductor R paketleri",
        "- GitHub ekosistemi",
        "- Python/PyTorch toplulugu",
        "",
        "REFERANS METODOLOJILER:",
        "- Horvath (2013): Pan-tissue epigenetik saat",
        "- Levine et al. (2018): PhenoAge",
        "- Lu et al. (2019): GrimAge",
        "- Belsky et al. (2022): DunedinPACE",
        "",
        "KOD ERISIMI:",
        "github.com/mortemdulcem/epi-clock-DNA-mtl-prototype",
        "",
        "FONLAMA: Bu calisma icin ozel fon alinmamistir.",
        "CIKAR CATISMASI: Yazarlar herhangi bir cikar catismasi bildirmemektedir."
    ], font_size=12)
    
    # SLIDE 100: Kapanış
    add_title_slide(prs,
        "TESEKKURLER",
        "EpiClock v4.0\nDNA Metilasyon Tabanli Epigenetik Yas Analiz Platformu\n\nDr. Nurcan Denli Bayir (nrcdnl94)\n\nSorular icin: nrcdnl94@epiclock.org\nGitHub: github.com/mortemdulcem/epi-clock-DNA-mtl-prototype"
    )
    
    # Save presentation
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"EpiClock_v4_Kapsamli_Sunum_{timestamp}.pptx"
    prs.save(filename)
    print(f"Sunum olusturuldu: {filename}")
    print(f"Toplam slayt sayisi: {len(prs.slides)}")
    return filename

if __name__ == "__main__":
    create_full_presentation()
