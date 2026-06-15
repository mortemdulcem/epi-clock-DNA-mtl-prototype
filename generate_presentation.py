"""
EpiClock v4.0 - Kapsamli Platform Sunumu
Tum Ozellikler ve Teorik Temeller
UNODC Kurumsal Tasarim

Author: Dr. Nurcan Denli Bayir (nrcdnl94)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime

UNODC_BLUE = RGBColor(0x00, 0x50, 0xA0)
UNODC_DARK = RGBColor(0x00, 0x33, 0x66)
UNODC_TURQUOISE = RGBColor(0x00, 0xA7, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = UNODC_DARK
    shape.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = UNODC_TURQUOISE
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = UNODC_TURQUOISE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_title, section_number):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = UNODC_BLUE
    shape.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"BOLUM {section_number}"
    p.font.size = Pt(20)
    p.font.color.rgb = UNODC_TURQUOISE
    p.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, two_column=False):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = UNODC_DARK
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if two_column and len(content_items) >= 2:
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(content_items[:len(content_items)//2]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = BLACK
            p.space_after = Pt(6)
        
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(content_items[len(content_items)//2:]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = BLACK
            p.space_after = Pt(6)
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = BLACK
            p.space_after = Pt(8)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = UNODC_DARK
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    cols = len(headers)
    table_rows = min(len(rows) + 1, 12)
    
    col_width = Inches(9) / cols
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
    
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = UNODC_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    for row_idx, row_data in enumerate(rows[:table_rows-1]):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = BLACK
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    
    return slide

def create_epiclock_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs, 
        "EpiClock v4.0",
        "DNA Metilasyonu Tabanli Epigenetik Yas Analiz Platformu\nBagimlilik Tespiti ve Adli Bilisim Uygulamalari\n\nDr. Nurcan Denli Bayir (nrcdnl94)"
    )
    
    add_section_slide(prs, "PLATFORM GENEL BAKIS", "1")
    
    add_content_slide(prs, "EpiClock v4.0 - Misyon ve Vizyon", [
        "MISYON: DNA metilasyonu yoluyla bagimlilik maddelerinin biyolojik etkilerini tespit etmek",
        "",
        "TEMEL YETENEKLER:",
        "   - 5 Major Epigenetik Saat (Horvath, Hannum, PhenoAge, GrimAge, DunedinPACE)",
        "   - 12 Doku-Spesifik Saat + Capraz-Doku Normalizasyonu",
        "   - Makine Ogrenmesi Ensemble Modeli (RF + XGBoost + ElasticNet)",
        "   - Graph Neural Network (GNN) Molekul Analizi",
        "   - 55+ Farmakolojik Reseptor Hedefi (IUPHAR/DrugBank/ChEMBL)",
        "   - UN/WHO/EMCDDA Regulatuvar Schedule Siniflandirmasi",
        "   - Blockchain Audit Trail ve Adli Bilisim Modulleri",
        "",
        "HEDEF KULLANICILAR:",
        "   - Klinik Arastirmacilar ve Hekimler",
        "   - Adli Tip Uzmanlari",
        "   - Bagimlilik Tedavi Merkezleri",
        "   - Regulatuvar Kurumlar (UNODC, WHO, EMCDDA)"
    ])
    
    add_content_slide(prs, "Sistem Mimarisi", [
        "FRONTEND KATMANI:",
        "   - Streamlit Web Uygulamasi (Port 5000)",
        "   - UNODC Kurumsal Tema (#0050A0, #003366, #00A7D8)",
        "   - Interaktif Plotly/Matplotlib Gorsellestirme",
        "   - Rol Tabanli Kullanici Arayuzu (Klinisyen/Arastirmaci/Adli)",
        "",
        "BACKEND KATMANI:",
        "   - Python 3.11 + NumPy/Pandas/SciPy",
        "   - PyTorch Graph Neural Network",
        "   - RDKit Molekul Islemleri",
        "   - Scikit-learn + XGBoost ML Pipeline",
        "",
        "VERITABANI KATMANI:",
        "   - PostgreSQL (Neon-backed)",
        "   - CpG Marker Veritabani",
        "   - Madde Panel Veritabani",
        "   - Audit Log ve Blockchain Kayitlari"
    ])
    
    add_section_slide(prs, "EPIGENETIK SAATLER", "2")
    
    add_content_slide(prs, "DNA Metilasyonu Teorisi", [
        "DNA METILASYONU NEDIR?",
        "   - Sitozin bazina metil grubu (-CH3) eklenmesi",
        "   - CpG dinukleotidlerinde yogunlasir (5'-CG-3')",
        "   - Epigenetik gen regülasyonunun temel mekanizmasi",
        "",
        "CPG ADALARI:",
        "   - Promotor bolgelerde yogunlasan CpG zengin bolgeler",
        "   - Genellikle metillenmemis (aktif transkripsiyon)",
        "   - Hipermetilasyon = Gen susturulmesi",
        "",
        "BIYOLOJIK YASLANMA:",
        "   - CpG metilasyon paternleri yasla degisir",
        "   - Belirli CpG siteleri yasla korelasyon gosterir",
        "   - 'Epigenetik saat' kavraminin temeli",
        "",
        "OLCUM YONTEMLERI:",
        "   - Illumina 450K/EPIC Array (485.000+ CpG)",
        "   - Bisulfit Sekanslama (WGBS, RRBS)",
        "   - Pyrosequencing (hedefli)"
    ])
    
    add_table_slide(prs, "5 Major Epigenetik Saat Karsilastirmasi",
        ["Saat", "Yil", "CpG Sayisi", "Doku", "Olcum", "Referans"],
        [
            ["Horvath", "2013", "353", "Pan-tissue", "DNAm Age", "Horvath, 2013"],
            ["Hannum", "2013", "71", "Kan", "DNAm Age", "Hannum et al., 2013"],
            ["PhenoAge", "2018", "513", "Kan", "Fenotipik Yas", "Levine et al., 2018"],
            ["GrimAge", "2019", "1030", "Kan", "Mortalite Riski", "Lu et al., 2019"],
            ["DunedinPACE", "2022", "173", "Kan", "Yaslanma Hizi", "Belsky et al., 2022"]
        ]
    )
    
    add_content_slide(prs, "Horvath Saati (2013)", [
        "TEORIK TEMEL:",
        "   - Ilk pan-tissue epigenetik saat",
        "   - 353 CpG sitesi ile kronolojik yas tahmini",
        "   - ElasticNet regresyon modeli",
        "",
        "MATEMATIKSEL FORMUL:",
        "   DNAm Age = anti.F(sum(w_i * beta_i) + intercept)",
        "   F = transformasyon fonksiyonu (adult/child ayirimi)",
        "",
        "OZELLIKLER:",
        "   - 51 farkli doku/hucre tipinde gecerli",
        "   - Median hata: 3.6 yil",
        "   - Korelasyon: r = 0.96",
        "",
        "KLINIK UYGULAMALAR:",
        "   - Biyolojik yas tahmini",
        "   - Yas hizlanmasi (Age Acceleration) hesabi",
        "   - Kanser, obezite, HIV ile iliskili calismalar"
    ])
    
    add_content_slide(prs, "PhenoAge ve GrimAge (2018-2019)", [
        "PHENOAGE (Levine, 2018):",
        "   - 513 CpG sitesi",
        "   - Fenotipik biyomarkerlerle egitilmis",
        "   - 9 klinik olcum: albumin, kreatinin, glukoz, CRP, lenfosit...",
        "   - Mortalite ve morbiditey daha iyi tahmin eder",
        "",
        "GRIMAGE (Lu, 2019):",
        "   - 1030 CpG sitesi",
        "   - DNAm-bazli plazma protein surrogatlari",
        "   - Sigara paket-yili, PAI-1, GDF-15, ADM, B2M...",
        "   - 'En iyi mortalite tahmincisi'",
        "",
        "DUNEDINPACE (Belsky, 2022):",
        "   - 173 CpG sitesi",
        "   - Yaslanma HIZI olcumu (pace of aging)",
        "   - Longitudinal Dunedin kohortu (45 yillik takip)",
        "   - Deger yorumu: 1.0 = normal, >1.0 = hizli yaslanma"
    ])
    
    add_section_slide(prs, "BAGIMLILIK VE EPIGENETIK YAS HIZLANMASI", "3")
    
    add_content_slide(prs, "Madde Bagimliligi ve Epigenetik Etkiler", [
        "TEMEL MEKANIZMALAR:",
        "   - Kronik madde kullanimi DNA metilasyon paternlerini degistirir",
        "   - Reward pathway genlerinde (OPRM1, DRD2, DAT) metilasyon degisiklikleri",
        "   - Stres yaniit genlerinde (NR3C1, FKBP5) epigenetik modifikasyonlar",
        "",
        "EPIGENETIK YAS HIZLANMASI (EAA):",
        "   - EAA = DNAm Age - Kronolojik Yas",
        "   - Pozitif EAA = Biyolojik olarak 'daha yasli'",
        "",
        "MADDE-SPESIFIK EAA DEGERLERI:",
        "   - Metamfetamin: +4.2 yil (95% CI: 3.5-4.9)",
        "   - Kokain: +3.1 yil (95% CI: 2.4-3.8)",
        "   - Opioidler: +2.9 yil (95% CI: 2.5-3.4)",
        "   - Alkol: +2.2 yil (95% CI: 1.8-2.6)",
        "   - Kannabis: +0.8 yil (95% CI: 0.3-1.4)",
        "",
        "KAYNAK: Rosen et al., 2018; Quach et al., 2017; Beach et al., 2015"
    ])
    
    add_content_slide(prs, "Madde Sinifina Gore CpG Degisiklikleri", [
        "OPIOIDLER (OPRM1, OPRD1, OPRK1):",
        "   - cg16419235 (OPRM1 promotor) - Hipometilasyon",
        "   - cg23663830 (POMC) - Hipermetilasyon",
        "   - Tolerans ve bagimlilik gelisimi ile iliskili",
        "",
        "STIMULANLAR (DAT, DRD2, COMT):",
        "   - cg06500161 (SLC6A3) - Metilasyon artisi",
        "   - cg14976642 (DRD2) - Promotor hipermetilasyonu",
        "   - Dopamin sinyal yolagi disregulasyonu",
        "",
        "ALKOL (ADH, ALDH, GABRA):",
        "   - cg21161138 (GABRA1) - Metilasyon degisimi",
        "   - cg09935388 (ADH1B) - Polimorfizm-metilasyon etkilesimi",
        "",
        "KANNABINOIDLER (CNR1, FAAH):",
        "   - cg12654770 (CNR1) - Kronik kullanimda hipometilasyon",
        "   - Endokannabinoid sistem disregulasyonu"
    ])
    
    add_section_slide(prs, "MAKINE OGRENMESI MODELLERI", "4")
    
    add_content_slide(prs, "Ensemble ML Arsitekturisi", [
        "MODEL BILESENLERI:",
        "   1. Random Forest (n_estimators=500, max_depth=20)",
        "   2. XGBoost (learning_rate=0.05, max_depth=6)",
        "   3. ElasticNet (alpha=0.5, l1_ratio=0.5)",
        "",
        "ENSEMBLE STRATEJISI:",
        "   - Weighted Average (agirlikli ortalama)",
        "   - RF: 0.35, XGBoost: 0.45, ElasticNet: 0.20",
        "   - Optimum agirliklar CV ile belirlenmis",
        "",
        "OZELLIK SETI:",
        "   - 10,542 CpG beta degeri",
        "   - Demografik ozellikler (yas, cinsiyet, etnisite)",
        "   - Klinik biyomarkerler (CRP, albumin, vb.)",
        "",
        "PERFORMANS METRIKLERI:",
        "   - RMSE: 2.8 yil",
        "   - MAE: 2.1 yil",
        "   - R-squared: 0.94",
        "   - 5-fold Cross-validation"
    ])
    
    add_section_slide(prs, "GRAPH NEURAL NETWORK MOLEKUL ANALIZI", "5")
    
    add_content_slide(prs, "GNN Arsitekturisi - MPNN + Attention", [
        "MODEL PARAMETRELERI:",
        "   - 4 Message Passing katmani",
        "   - 256 gizli boyut",
        "   - Multi-head Attention (8 head)",
        "   - ~1.2 milyon parametre",
        "",
        "ATOM OZELLIKLERI (146 boyut):",
        "   - Atom numarasi (one-hot encoding)",
        "   - Derece, formal yuk, hibridizasyon",
        "   - Aromatiklik, halka uyeligi",
        "   - Implicit hidrojenler, radikallik",
        "",
        "BAG OZELLIKLERI (12 boyut):",
        "   - Bag tipi (tek, cift, uc, aromatik)",
        "   - Konjugasyon, halka icinde olma",
        "   - Stereo konfigurasyonu",
        "",
        "CIKIS KATMANLARI:",
        "   - Bagimlilik potansiyeli + Belirsizlik",
        "   - Toksisite siniflandirmasi",
        "   - Metabolizma yukumlulugu",
        "   - CYP inhibisyonu (1A2, 2C9, 2C19, 2D6, 3A4)"
    ])
    
    add_content_slide(prs, "Molekuler Graf Gosterimi", [
        "SMILES -> GRAF DONUSUMU:",
        "   - RDKit ile molekul okuma",
        "   - Her atom = bir dugum (node)",
        "   - Her bag = bir kenar (edge)",
        "   - Yonlu graf (directed graph)",
        "",
        "MESSAGE PASSING MEKANIZMASI:",
        "   m_ij = MLP([h_i || h_j || e_ij])",
        "   h_i' = GRU(h_i, sum(m_ji))",
        "",
        "GLOBAL HAVUZLAMA:",
        "   - Attention-weighted sum",
        "   - Max pooling + Mean pooling",
        "   - Molekul seviyesi temsil",
        "",
        "UNCERTAINTY QUANTIFICATION:",
        "   - Monte Carlo Dropout",
        "   - 50 forward pass ile belirsizlik",
        "   - Epistemic + Aleatoric ayirimi"
    ])
    
    add_section_slide(prs, "GELISMIS OZELLIK MUHENDISLIGI", "6")
    
    add_table_slide(prs, "55+ Reseptor Hedefi - Kategoriler",
        ["Kategori", "Reseptorler", "Bagimlilik Agirligi"],
        [
            ["Opioid", "MOR, DOR, KOR, NOP", "%25-95"],
            ["Dopaminerjik", "DAT, D1, D2, D3, D4, D5", "%50-92"],
            ["Serotonerjik", "SERT, 5-HT1A/1B/2A/2C/3", "%30-60"],
            ["GABAerjik", "GABA-A (alpha1-5, gamma2, delta), GABA-B", "%45-75"],
            ["Glutamaterjik", "NMDA NR1/NR2B, AMPA, mGluR2/5", "%40-60"],
            ["Kannabinoid", "CB1, CB2", "%25-55"],
            ["Kolinerjik", "nAChR (a4b2, a7, a3b4), mAChR (M1, M2)", "%30-85"],
            ["Adrenerjik", "NET, Alpha2A, Beta1", "%35-65"],
            ["Stres/Neuropeptid", "CRF1/2, OX1/2, NK1, NPY-Y1", "%40-65"],
            ["Sigma", "Sigma-1, Sigma-2", "%35-50"],
            ["Enzimler", "MAO-A/B, COMT, FAAH", "%40-45"]
        ]
    )
    
    add_content_slide(prs, "Kimyasal Ozellik Ekstraksiyon", [
        "MORGAN/ECFP FINGERPRINTS:",
        "   - 2048-bit dairesel fingerprint",
        "   - Radius = 2 (ECFP4 esdegeri)",
        "   - Yapısal benzerlik hesabi",
        "",
        "FIZIKOKIMYASAL DESCRIPTORLER:",
        "   - Molekul agirligi (MW)",
        "   - cLogP (lipofili)",
        "   - TPSA (polar yuzey alani)",
        "   - HBD/HBA (hidrojen bag verici/alici)",
        "   - Rotatable bond sayisi",
        "   - Aromatik halka sayisi",
        "",
        "PKA TAHMINLERI:",
        "   - Asidik/bazik gruplar icin pKa",
        "   - Iyonizasyon durumu tahmini",
        "",
        "LIPINSKI/VEBER KURALLARI:",
        "   - MW < 500, LogP < 5, HBD <= 5, HBA <= 10",
        "   - TPSA < 140, Rotatable bonds <= 10",
        "   - Ihlal sayisi hesabi"
    ])
    
    add_content_slide(prs, "Farmakokinetik Profil Analizi", [
        "BBB GECIRGENLIK:",
        "   - LogP ve TPSA bazli tahmin",
        "   - Optimal: LogP 2-4, TPSA < 70",
        "   - Yuksek BBB = Yuksek CNS etkisi",
        "",
        "PLAZMA PROTEIN BAGLAMA:",
        "   - Albumin ve AAG baglama tahmini",
        "   - Serbest fraksiyon hesabi",
        "",
        "YARI OMUR (t1/2):",
        "   - Hepatik klerens tahmini",
        "   - Kisa t1/2 = Sik doz tekrari = Yuksek bagimlilik",
        "",
        "BIYOYARARLANIM:",
        "   - Oral absorpsiyon tahmini",
        "   - Ilk gecis metabolizmasi",
        "",
        "CYP ENZIM INHIBISYONU:",
        "   - CYP1A2, CYP2C9, CYP2C19, CYP2D6, CYP3A4",
        "   - Ilac-ilac etkilesim riski"
    ])
    
    add_section_slide(prs, "REGULATUVAR CERCEVE", "7")
    
    add_table_slide(prs, "UN/WHO/EMCDDA Schedule Siniflandirmasi",
        ["Schedule", "Tanim", "Istismar Skoru", "Ornekler"],
        [
            ["I", "Yuksek istismar, tibbi kullanim yok", "%95", "Heroin, LSD, MDMA"],
            ["II", "Yuksek istismar, sinirli tibbi", "%85", "Morfin, Fentanil, Kokain"],
            ["III", "Orta istismar, tibbi kullanim var", "%65", "Buprenorfin, Ketamin"],
            ["IV", "Dusuk-orta istismar", "%45", "Benzodiazepinler, Tramadol"],
            ["V", "Dusuk istismar", "%25", "Pregabalin, Dusuk doz kodein"],
            ["Unscheduled", "Kontrolsuz", "%10", "Kafein (Nikotin/Alkol ayri)"]
        ]
    )
    
    add_content_slide(prs, "Istismar Potansiyeli Hesaplama", [
        "WEAK-SUPERVISION YAKLASIMI:",
        "   - Validated substances: Literatur referansli degerler",
        "   - Predicted substances: Reseptor + PK faktörleri",
        "",
        "HESAPLAMA FORMULU:",
        "   Abuse_Score = w1*Receptor_Score + w2*PK_Score + w3*Schedule_Score",
        "",
        "RESEPTOR KATKISI:",
        "   - MOR aktivasyonu: +0.95 (en yuksek)",
        "   - DAT inhibisyonu: +0.92",
        "   - GABA-A modulasyonu: +0.75",
        "   - Toplamda 55+ reseptor agirligi",
        "",
        "FARMAKOKINETIK CARPANLAR:",
        "   - Hizli etki baslangici: x1.5",
        "   - Yuksek BBB gecisi: x1.3",
        "   - Kisa yari omur: x1.2",
        "   - Yuksek biyoyararlanim: x1.1",
        "",
        "GUVEN ARALIGI: Bootstrap ile %95 CI hesabi"
    ])
    
    add_section_slide(prs, "ADLI BILISIM MODULLERI", "8")
    
    add_content_slide(prs, "Blockchain Audit Trail", [
        "BLOCKCHAIN OZELLIKLERI:",
        "   - SHA-256 hash zinciri",
        "   - Degistirilemez kayit (immutable)",
        "   - Zaman damgali islemler",
        "",
        "AUDIT KAYITLARI:",
        "   - Ornek girisi ve kimlik dogrulama",
        "   - Analiz baslangic/bitis zamanlari",
        "   - Sonuc uretimi ve onaylama",
        "   - Kullanici erisim loglari",
        "",
        "TAMPER DETECTION:",
        "   - Hash zinciri butunluk kontrolu",
        "   - Herhangi bir degisiklik tespit edilir",
        "   - Adli delil kabul edilebilirligi",
        "",
        "CHAIN OF CUSTODY:",
        "   - Ornek transferi takibi",
        "   - Sorumluluk zinciri dokumantasyonu",
        "   - Dijital imza destegi"
    ])
    
    add_content_slide(prs, "Postmortem Validasyon", [
        "PMI (POSTMORTEM INTERVAL) DUZELTMESI:",
        "   - DNA metilasyonu olum sonrasi degisir",
        "   - PMI-bagimlı beta deger ayarlamasi",
        "   - Doku tipine gore duzeltme faktorleri",
        "",
        "DOKU BUTUNLUK DEGERLENDIRMESI:",
        "   - DNA kalitesi skoru",
        "   - Degradasyon markerleri",
        "   - Guvenirlirlik skoru hesabi",
        "",
        "ADLI TIP UYGULAMALARI:",
        "   - Olum ani madde kullanim tespiti",
        "   - Kronik bagimlilik gecmisi",
        "   - Cause of death destekleme",
        "",
        "RAPORLAMA:",
        "   - PDF klinik rapor uretimi",
        "   - Adli standartlara uygun format",
        "   - Yorum ve oneriler"
    ])
    
    add_section_slide(prs, "VERITABANI VE ENTEGRASYONLAR", "9")
    
    add_content_slide(prs, "Veritabani Semaşi", [
        "POSTGRESQL TABLOLARI:",
        "",
        "cpg_markers:",
        "   - marker_id, chromosome, position, gene, island",
        "   - clock_membership (Horvath, Hannum, etc.)",
        "   - methylation_mean, std, range",
        "",
        "substance_panels:",
        "   - substance_id, name, schedule, class",
        "   - abuse_potential, references",
        "   - receptor_targets, metabolism_pathway",
        "",
        "analysis_results:",
        "   - sample_id, patient_id, analysis_date",
        "   - clock_results (JSON), eaa_values",
        "   - ml_predictions, gnn_results",
        "",
        "audit_logs:",
        "   - log_id, timestamp, user_id, action",
        "   - hash_chain, previous_hash"
    ])
    
    add_content_slide(prs, "Dis Veritabani Entegrasyonlari", [
        "GENOMIK VERITABANLARI:",
        "   - 1000 Genomes Project",
        "   - gnomAD (Genome Aggregation Database)",
        "   - UK Biobank",
        "   - TOPMed",
        "",
        "EPIGENETIK VERITABANLARI:",
        "   - EWAS Catalog",
        "   - GEO (Gene Expression Omnibus)",
        "   - ENCODE (Encyclopedia of DNA Elements)",
        "",
        "FARMAKOLOJIK VERITABANLARI:",
        "   - IUPHAR/BPS Guide to Pharmacology (3,112 hedef)",
        "   - DrugBank 6.0 (4,563 onayii ilac)",
        "   - ChEMBL (13,503 ligand)",
        "   - PharmGKB (Farmakogenomik)",
        "",
        "REGULATUVAR KAYNAKLAR:",
        "   - UN 1961/1971 Conventions",
        "   - WHO Expert Committee raporlari",
        "   - EMCDDA Early Warning System"
    ])
    
    add_section_slide(prs, "KLINIK UYGULAMALAR", "10")
    
    add_content_slide(prs, "Klinik Kullanim Senaryolari", [
        "BAGIMLILIK TEDAVI MERKEZLERI:",
        "   - Tedavi oncesi biyolojik yas degerlendirmesi",
        "   - Tedavi etkinligi takibi (EAA degisimi)",
        "   - Relaps riski tahmini",
        "",
        "ADLI TIP:",
        "   - Madde kullanim gecmisi tespiti",
        "   - Postmortem bagimlilik analizi",
        "   - Delil destekleme",
        "",
        "KLINIK ARASTIRMA:",
        "   - Madde-spesifik epigenetik etkiler",
        "   - Longitudinal kohort calismalari",
        "   - Biomarker kesfii",
        "",
        "KORUYUCU SAGLIK:",
        "   - Risk degerlendirmesi",
        "   - Erken mudahale planlama",
        "   - Genetik yatkinlik entegrasyonu"
    ])
    
    add_content_slide(prs, "Karar Destek Sistemi", [
        "RISK SKORLAMASI:",
        "   - Dusuk Risk (EAA < 1 yil): Takip",
        "   - Orta Risk (EAA 1-3 yil): Yakin izlem",
        "   - Yuksek Risk (EAA > 3 yil): Mudahale",
        "",
        "TEDAVI ONERILERI:",
        "   - Madde-spesifik tedavi protokolleri",
        "   - Farmakogenetik uyumluluk",
        "   - Yasam tarzı modifikasyonlari",
        "",
        "TAKIP PROTOKOLU:",
        "   - 3-6 aylik DNA metilasyon kontrolu",
        "   - EAA trendii izleme",
        "   - Tedavi yanit degerlendirmesi",
        "",
        "RAPORLAMA:",
        "   - Klinisyen-dostu ozet rapor",
        "   - Detayli teknik rapor",
        "   - PDF export (ReportLab)"
    ])
    
    add_section_slide(prs, "TEKNIK DOKUMANTASYON", "11")
    
    add_table_slide(prs, "Python Kutuphaneleri",
        ["Kutuphane", "Versiyon", "Kullanim"],
        [
            ["streamlit", "latest", "Web arayuzu"],
            ["pandas", "latest", "Veri islemleri"],
            ["numpy", "latest", "Numerik hesaplamalar"],
            ["scikit-learn", "latest", "ML modelleri"],
            ["xgboost", "latest", "Gradient boosting"],
            ["torch", "latest", "GNN/Deep learning"],
            ["rdkit", "latest", "Molekul islemleri"],
            ["plotly", "latest", "Interaktif grafikler"],
            ["matplotlib/seaborn", "latest", "Statik grafikler"],
            ["psycopg2", "latest", "PostgreSQL baglantisi"],
            ["reportlab", "latest", "PDF uretimi"]
        ]
    )
    
    add_content_slide(prs, "API ve Entegrasyon", [
        "REST API ENDPOINTLERI:",
        "   - POST /api/analyze - DNA metilasyon analizi",
        "   - POST /api/molecule - Molekul GNN analizi",
        "   - GET /api/results/{id} - Sonuc sorgulama",
        "   - GET /api/audit/{sample_id} - Audit log",
        "",
        "VERI FORMATLARI:",
        "   - Input: CSV, Excel, IDAT dosyalari",
        "   - Output: JSON, PDF, Excel rapor",
        "   - SMILES/SDF molekul formatlari",
        "",
        "GUVENLIK:",
        "   - Session-based authentication",
        "   - HTTPS zorunlu",
        "   - Hassas veri sifreleme",
        "",
        "INTEROPERABILITY:",
        "   - HL7 FHIR uyumlu",
        "   - LOINC kodlari destegi"
    ])
    
    add_section_slide(prs, "SONUC VE GELECEK YONELIMLER", "12")
    
    add_content_slide(prs, "Platform Ozeti", [
        "EPICLOCK v4.0 YETENEKLERI:",
        "",
        "   [+] 5 Major Epigenetik Saat + 12 Doku-Spesifik Saat",
        "   [+] Ensemble ML: RF + XGBoost + ElasticNet (R2=0.94)",
        "   [+] PyTorch GNN: MPNN + Attention (1.2M parametre)",
        "   [+] 55+ Farmakolojik Reseptor Hedefi",
        "   [+] UN/WHO/EMCDDA Regulatuvar Uyum",
        "   [+] 40+ Valide Madde Veritabani",
        "   [+] Blockchain Audit Trail",
        "   [+] Adli Bilisim Modulleri",
        "   [+] PostgreSQL Veritabani Entegrasyonu",
        "   [+] Interaktif Web Arayuzu",
        "",
        "ONEMLI NOT:",
        "   Bu bir PROTOTIP'tir. Gercek klinik/adli kullanim icin",
        "   UCSD lisanslari ve valide koefisyenler gerekmektedir."
    ])
    
    add_content_slide(prs, "Gelecek Gelistirmeler", [
        "KISA VADELI (6 ay):",
        "   - Gercek epigenetik saat koefisyenlerinin lisanslanmasi",
        "   - Daha buyuk valide madde veritabani",
        "   - Multi-omics entegrasyonu (transkriptomik, proteomik)",
        "",
        "ORTA VADELI (1 yil):",
        "   - Prospektif klinik validasyon calismalari",
        "   - Mobile uyumlu arayuz",
        "   - Federated learning ile gizlilik korumali ogrenme",
        "",
        "UZUN VADELI (2+ yil):",
        "   - Regulatory approval (CE-IVD, FDA clearance)",
        "   - Uluslararasi cok merkezli calismalar",
        "   - Yapay zeka destekli tedavi optimizasyonu",
        "",
        "ILETISIM:",
        "   Dr. Nurcan Denli Bayir (nrcdnl94)",
        "   EpiClock Gelistirme Ekibi"
    ])
    
    add_title_slide(prs, 
        "TESEKKURLER",
        "Sorular ve Tartisma\n\nEpiClock v4.0\nDr. Nurcan Denli Bayir"
    )
    
    filename = f"EpiClock_v4_Sunum_{datetime.now().strftime('%Y%m%d')}.pptx"
    prs.save(filename)
    print(f"Sunum olusturuldu: {filename}")
    return filename

if __name__ == "__main__":
    create_epiclock_presentation()
